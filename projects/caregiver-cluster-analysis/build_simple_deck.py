"""ESD Lab plain-language deck: 10 slides, original figures, simple words.

Same verified numbers as the full deck. Nothing here is a new analysis.
Hard rule enforced at the end of this script: no em dashes and no en dashes
anywhere in the slide text.
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import pandas as pd
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from esd_deck_lib import *

OUTDIR = os.path.dirname(OUT)
F = lambda n: f"{OUT}/{n}"

# ── ground truth, re-read and asserted (same anchors as the full deck) ────────
t15  = pd.read_csv(F("table_15_rule_false_positive_rates.csv"))
t15b = pd.read_csv(F("table_15b_rule_counts_by_project.csv"))
t16  = pd.read_csv(F("table_16_trust_tier_counts.csv"))
t34  = pd.read_csv(F("table_34_branching_logic_audit.csv"))
assert int(t16.n.sum()) == 1956
assert tuple(t15b[t15b.rule == "R8"].iloc[0][["all_n", "clean_4797_n", "dirty_4581_n"]]) == (44, 0, 44)
assert float(t15[t15.rule == "R8"].false_positive_pct.iloc[0]) == 0.0
assert len(t34) == 44 and set(t34.source_project) == {"dirty_4581"}
TIER = {p: {int(r.tier): (int(r.n), float(r.pct_within_project))
            for _, r in t16[t16.source_project == p].iterrows()}
        for p in ("clean_4797", "dirty_4581")}
print("[anchors] verified against source CSVs")

prs = new_deck()

# ═════════════════════════════════════════════════════════════════════════════
# LAYOUT HELPERS FOR THE PLAIN-LANGUAGE DECK
# ═════════════════════════════════════════════════════════════════════════════
def head2(s, title, sub):
    txt(s, M, 0.44, 11.15, 0.55, title, size=27, color=DISCOVERY, font=FH, bold=True,
        space_after=0, line=0.93, char_space=-0.2)
    txt(s, M, 1.12, 11.55, 0.62, sub, size=13, color=JET, font=FB, space_after=0, line=1.04)
    lab_logo(s, SW - M - 0.62, 0.44, 0.29)


from PIL import ImageFont

for _font_path in [
    f"{SK}/assets/fonts/LibreFranklin-Medium.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
]:
    try:
        _FT = ImageFont.truetype(_font_path, 200)
        break
    except Exception:
        continue
else:
    _FT = ImageFont.load_default()

def n_lines(text, width_in, pt):
    """Exact wrapped-line count, measured with the real Libre Franklin metrics."""
    limit = width_in * 72.0 * (200.0 / pt)      # available width in 200pt-font units
    lines, cur = 1, ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if _FT.getlength(trial) <= limit or not cur:
            cur = trial
        else:
            lines += 1; cur = word
    return lines


def read_card(s, x, y, w, h, points, stacked=True, title="How to read this"):
    plate(s, x, y, w, h, fill=COOLBLUE, radius=0.11)
    txt(s, x + 0.26, y + 0.19, w - 0.52, 0.22, title, size=9.5, color=ORANGE, font=FH,
        bold=True, caps=True, space_after=0, char_space=0.6)
    if stacked:
        yy = y + 0.52
        for i, p in enumerate(points):
            cir = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.26), Inches(yy + 0.015),
                                     Inches(0.24), Inches(0.24))
            cir.fill.solid(); cir.fill.fore_color.rgb = DISCOVERY
            cir.line.fill.background(); cir.shadow.inherit = False
            txt(s, x + 0.26, yy + 0.045, 0.24, 0.20, str(i + 1), size=9, color=COOLWHITE,
                font=FH, bold=True, align=PP_ALIGN.CENTER, space_after=0)
            txt(s, x + 0.60, yy, w - 0.88, 0.9, p, size=11, color=JET, font=FB,
                space_after=0, line=1.02)
            yy += 0.29 + 0.1585 * n_lines(p, w - 0.88, 11)
    else:
        cw = (w - 0.52 - 0.34 * (len(points) - 1)) / len(points)
        for i, p in enumerate(points):
            xx = x + 0.26 + i * (cw + 0.34)
            cir = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xx), Inches(y + 0.50),
                                     Inches(0.24), Inches(0.24))
            cir.fill.solid(); cir.fill.fore_color.rgb = DISCOVERY
            cir.line.fill.background(); cir.shadow.inherit = False
            txt(s, xx, y + 0.53, 0.24, 0.20, str(i + 1), size=9, color=COOLWHITE, font=FH,
                bold=True, align=PP_ALIGN.CENTER, space_after=0)
            txt(s, xx + 0.34, y + 0.50, cw - 0.34, h - 0.66, p, size=10.5, color=JET,
                font=FB, space_after=0, line=1.03)


def caution(s, text, color=None):
    plate(s, M, 6.88, SW - 2 * M, 0.40, fill=COOLWHITE, radius=0.18)
    txt(s, M + 0.24, 6.965, SW - 2 * M - 0.48, 0.26, text, size=10.5,
        color=color or RED, font=FH, bold=True, space_after=0)


def simple_fig_slide(prs, title, sub, fig, points, src, notes, caut=None):
    s = blank(prs)
    head2(s, title, sub)
    a = aspect(fig)
    fy, fh, cy = (1.80, 3.34, 5.26) if caut else (1.86, 3.62, 5.62)
    if a >= 1.9:                                   # wide figure, card underneath
        figure(s, fig, M, fy, SW - 2 * M, fh)
        read_card(s, M, cy, SW - 2 * M, 1.16, points, stacked=False)
    else:                                          # tall figure, card beside it
        hh = 4.42 if caut else 4.92
        figure(s, fig, M, fy, 7.62, hh)
        read_card(s, 8.42, fy, 4.31, hh, points, stacked=True)
    s.notes_slide.notes_text_frame.text = notes
    if caut:
        txt(s, M, 6.52, SW - 2 * M, 0.22, "Source: " + src, size=8, color=GREY,
            font=FB, space_after=0)
        caution(s, caut)
    else:
        source(s, src)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# 1. TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = blank(prs, bg=DISCOVERY)
if os.path.exists(PATTERN_BAND):
    s.shapes.add_picture(PATTERN_BAND, Inches(-0.6), Inches(5.55), width=Inches(14.5), height=Inches(2.4))
plate(s, M, 0.46, 5.95, 0.92, fill=WHITE, radius=0.28)
logo_pair(s, M + 0.32, 0.66, 0.46)
txt(s, M, 2.20, 11.6, 1.6, "Caregiver Screening Attitudes\nand Data Quality",
    size=42, color=COOLWHITE, font=FH, bold=True, line=0.94, char_space=-0.3, space_after=0)
txt(s, M, 3.92, 11.6, 0.36, "A plain language walk through of what we found",
    size=16, color=COOLBLUE, font=FB, space_after=0)
txt(s, M, 4.42, 11.6, 0.60,
    ["Early Social Development Lab, University of South Carolina",
     "Based on 1,956 survey records. Same numbers as the full technical deck, fewer of them on each slide."],
    size=10.5, color=SCIENCE, font=FB, space_after=3)
s.notes_slide.notes_text_frame.text = (
    "I put this version together so we can talk through the findings without anyone needing to hold ten numbers in their head at once. "
    "Every figure here also appears in the long deck, and every number is the same. I have just taken the scaffolding off.\n\n"
    "There are two stories running side by side. The first is about caregivers and how they feel about infant autism screening. "
    "The second is about how much of the incoming survey data we can actually trust. I will do the caregiver story first, then the data "
    "quality story, and then finish with what we are willing to say out loud and what we are not.\n\n"
    "If anything is unclear as I go, please stop me. That is the whole point of this version."
)

# ═════════════════════════════════════════════════════════════════════════════
# 2. WHAT WE DID
# ═════════════════════════════════════════════════════════════════════════════
s = blank(prs)
head2(s, "What we did, in one slide",
      "We asked caregivers how they feel about infant autism screening. Then, separately, we checked whether those feelings line up with what they say they would actually do.")
# left panel: goes in
plate(s, M, 2.00, 5.40, 3.05, fill=COOLBLUE, radius=0.14)
txt(s, M + 0.30, 2.20, 4.80, 0.26, "What builds the groups", size=10, color=ORANGE,
    font=FH, bold=True, caps=True, space_after=0, char_space=0.6)
txt(s, M + 0.30, 2.52, 4.80, 0.40, "10 attitude questions only", size=17, color=DISCOVERY,
    font=FH, bold=True, space_after=0)
txt(s, M + 0.30, 3.00, 4.80, 1.85,
    "Trust in clinicians, comfort with a blood test, whether screening feels ethical, "
    "how much distress it would cause, how useful it seems, tolerance for a wrong answer, "
    "openness to equipment, feasibility, low burden options, and how unpleasant it feels.",
    size=10.5, color=JET, font=FB, space_after=0, line=1.05)
# arrow
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.16), Inches(3.24), Inches(0.92), Inches(0.58))
ar.fill.solid(); ar.fill.fore_color.rgb = DISCOVERY; ar.line.fill.background(); ar.shadow.inherit = False
# right panel: held out
plate(s, 7.33, 2.00, 5.40, 3.05, fill=COOLWHITE, radius=0.14)
txt(s, 7.63, 2.20, 4.80, 0.26, "Kept out, checked afterwards", size=10, color=ORANGE,
    font=FH, bold=True, caps=True, space_after=0, char_space=0.6)
txt(s, 7.63, 2.52, 4.80, 0.40, "4 things we never let in", size=17, color=DISCOVERY,
    font=FH, bold=True, space_after=0)
for i, (lab, dsc) in enumerate([
        ("Screening intent", "Would you get your baby screened?"),
        ("Autism knowledge", "A 7 question quiz with right answers"),
        ("Personal values", "What matters to you in life"),
        ("Demographics", "Education, work, area, family setup")]):
    yy = 3.02 + i * 0.47
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.63), Inches(yy + 0.055), Inches(0.13), Inches(0.13))
    dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE; dot.line.fill.background(); dot.shadow.inherit = False
    txt(s, 7.88, yy, 2.05, 0.24, lab, size=11.5, color=JET, font=FH, bold=True, space_after=0)
    txt(s, 9.95, yy + 0.015, 2.55, 0.24, dsc, size=10, color=GREY, font=FB, space_after=0)
read_card(s, M, 5.30, SW - 2 * M, 1.16, [
    "1,956 records came in. 131 caregivers answered enough attitude questions to be grouped. 127 of those also answered the screening question.",
    "The screening answer never went into the grouping, so when we compare the groups on it later we are not arguing in a circle.",
    "Everything on the right is a check we run after the groups already exist. It cannot bend the groups to fit a story.",
], stacked=False)
source(s, "table_5b_cohort_flow.csv, table_16_trust_tier_counts.csv, config.yaml")
s.notes_slide.notes_text_frame.text = (
    "This is the one design idea that everything else rests on, so it is worth thirty seconds even with a friendly audience.\n\n"
    "We built the caregiver groups using only the ten attitude questions on the left. The four things on the right were deliberately locked "
    "out of that process. We only looked at them once the groups already existed and could not be changed.\n\n"
    "Why that matters: if I had let the screening answer help build the groups, and then announced that the groups differ on screening, "
    "I would have proved nothing at all. Keeping it out is what makes the comparison on the next few slides worth anything.\n\n"
    "On the numbers, 1,956 records came in across the two projects, but only 131 caregivers answered enough of the attitude block to be "
    "placed in a group, and 127 of those also gave a usable screening answer. Those are the people the caregiver story is about."
)

# ═════════════════════════════════════════════════════════════════════════════
# 3. TWO GROUPS
# ═════════════════════════════════════════════════════════════════════════════
simple_fig_slide(prs,
    "Caregivers fall into two broad groups",
    "One group is warm about screening across the board. The other is supportive too, but with conditions attached.",
    F("figure_2_primary_cluster_profiles.png"),
    ["Each row is one attitude. Further right means more supportive. Blue circles are the warmer group, orange squares are the conditional group.",
     "The wider the gap between the two markers, the more that attitude separates the groups. The three widest gaps are all about trust and ethics.",
     "Look at the bottom row. Both groups sit at the far right, so neither of them finds the screen itself upsetting. The disagreement is about the conditions around it."],
    "figure_2_primary_cluster_profiles.png, table_4_cluster_defining_profiles.csv",
    "I want to be careful with the labels here. The second group is not against screening. If you look at where their orange markers sit, "
    "most of them are above the midpoint. That is why we call them conditional rather than low acceptability.\n\n"
    "What actually splits the two groups is trust in clinicians, willingness to accept a blood test, and whether screening feels ethically "
    "right to them. Those three sit at the top because they have the widest gaps, all around 25 points on a 100 point scale.\n\n"
    "The bottom row is the one I find most useful. Both groups agree that the screen itself is not distressing. So this is not squeamishness "
    "about the procedure. It is about whether they trust the people and the process around it, which is a much more actionable thing for us to work on.\n\n"
    "One honest caveat if someone asks. These are two ends of a spread rather than two clean types of person. The groups are a useful summary, "
    "not a discovery of two different kinds of caregiver."
)

# ═════════════════════════════════════════════════════════════════════════════
# 4. SCREENING INTENT
# ═════════════════════════════════════════════════════════════════════════════
simple_fig_slide(prs,
    "The two groups answer the screening question very differently",
    "About 7 in 10 caregivers in the warmer group said they would definitely screen their baby. In the conditional group it was closer to 3 in 10.",
    F("figure_3_screening_outcome.png"),
    ["Left panel: the share who said definitely yes. The vertical bar shows how sure we are about that share.",
     "The two bars do not overlap, which is the simple test for whether a difference is real or just noise. Here it is real.",
     "Right panel: the full mix of answers. Notice the dark blue block is much bigger for the warmer group."],
    "figure_3_screening_outcome.png, table_5_screening_outcome.csv",
    "This is the headline, and it is a big gap. Seventy percent against twenty eight percent, so roughly forty two points apart.\n\n"
    "The way I would put it in a sentence for a non technical reader is this. If you know which of the two attitude groups a caregiver falls "
    "into, you can predict their screening answer far better than chance, and we worked that out without ever letting the screening answer "
    "influence the grouping.\n\n"
    "The thing I want to head off is the causal reading. Nobody should walk away thinking that changing a caregiver's attitude will change "
    "their screening decision. We asked both questions in the same sitting, so we have no way to say which came first, or whether something "
    "else drives both. It is a strong link, and links are useful for targeting outreach, but it is not a lever we have proven.\n\n"
    "If someone worries the result depends on where we drew the line, we also checked the looser version. Counting probably yes as well, "
    "it is ninety percent against sixty percent. Same ordering, so it is not an artefact of the cut point.",
    caut="Careful: this is a link, not a cause. We asked about attitudes and screening in the same sitting, so we cannot say which one drives the other."
)

# ═════════════════════════════════════════════════════════════════════════════
# 5. ADJUSTED
# ═════════════════════════════════════════════════════════════════════════════
simple_fig_slide(prs,
    "It is not simply families who already live with autism",
    "Having an autistic child at home does make a caregiver more likely to screen. It does not explain the group difference away.",
    F("figure_18_logistic_forest.png"),
    ["Each dot is an odds ratio. Anything to the right of the dashed line at 1 means more likely to say definitely yes.",
     "Bottom row: the attitude group is worth about 8.7 times the odds, even after we account for autism in the home.",
     "Middle row: autism at home matters on its own, about 4.3 times the odds.",
     "Top row asks whether the two effects feed off each other. The line crosses 1, so we cannot tell. They look like they just add up."],
    "figure_18_logistic_forest.png, table_22_logistic_models.csv, table_11_asd_family_context.csv",
    "The obvious challenge to the previous slide is that families already living with autism are both more accepting and more likely to "
    "screen, and that this is the whole story. So we tested it head on.\n\n"
    "Autism at home does matter, about four times the odds, and that is a real finding in its own right. But the attitude group effect barely "
    "moves once we account for it. It stays around eight and a half times the odds. So the two things are working alongside each other rather "
    "than one hiding behind the other.\n\n"
    "The simplest way to see it is the split by household. Among families with an autistic child at home it is eighty eight percent against "
    "thirty eight percent. Among families without, it is fifty seven percent against seventeen percent. Same direction, same rough size, both clear.\n\n"
    "The top row is the one I would not over read. It asks whether being in the warmer group matters more when you already have autism in the "
    "family. The answer is that we cannot tell. With this sample size, that line would have to be much narrower before it told us anything."
)

# ═════════════════════════════════════════════════════════════════════════════
# 6. WHAT DOES NOT SEPARATE
# ═════════════════════════════════════════════════════════════════════════════
simple_fig_slide(prs,
    "Some things you would expect to matter, do not",
    "Autism knowledge, personal values and demographics all look about the same in both groups. That is a finding worth stating carefully.",
    F("figure_4_autism_knowledge.png"),
    ["Left panel: average score on a 7 question autism quiz. The vertical bars overlap, so we cannot call this a real difference.",
     "Right panel: the same comparison, question by question. Both groups are strong on vaccines and Tylenol, both weak on genetics.",
     "The same picture holds for values and for demographics. Nothing survives once we account for how many things we tested."],
    "figure_4_autism_knowledge.png, figure_5_values_characterization.png, figure_6_demographic_characterization.png, table_29_power_precision.csv",
    "I like this slide because it keeps us honest. The two groups differ enormously on attitudes and on screening intent, but they are not "
    "obviously different kinds of people.\n\n"
    "Knowledge goes the way you would expect, about 3.9 correct against 3.4 out of seven, but the uncertainty bars overlap and I would not "
    "build anything on it. Values looked more interesting at first, with self direction showing the largest raw difference, but we tested ten "
    "values and once you account for that it stops being convincing. Demographics show nothing at all.\n\n"
    "Now the part I would put in bold. With forty seven people in one group and eighty in the other, we could only reliably spot a difference "
    "of about twenty five percentage points. Anything smaller than that is simply invisible to us. So the correct sentence is that we did not "
    "find a difference, not that there is no difference. Please do not let anyone write that the two groups are demographically identical.",
    caut="We did not find a difference. That is not the same as there being none. At this sample size, gaps under about 25 points would be invisible to us."
)

# ═════════════════════════════════════════════════════════════════════════════
# 7. DATA TRUST
# ═════════════════════════════════════════════════════════════════════════════
s = blank(prs)
head2(s, "Now the second question: how much of this data can we trust?",
      "One project was collected under supervision. The other was flooded with low quality submissions. We score every record against the same fixed set of rules.")
cd = CategoryChartData()
cd.categories = ["Clean project\n(177 records)", "Legacy project\n(1,779 records)"]
for tier, label in [(1, "Tier 1  Something impossible in it"), (2, "Tier 2  Looks very doubtful"),
                    (3, "Tier 3  Not sure either way"), (4, "Tier 4  Passed every check")]:
    cd.add_series(label, (round(TIER["clean_4797"][tier][1], 1), round(TIER["dirty_4581"][tier][1], 1)))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, Inches(M), Inches(1.86), Inches(7.62), Inches(3.70), cd)
ch = gf.chart
ch.has_title = False
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(9); ch.legend.font.name = FB
for ser, col in zip(ch.plots[0].series, [RED, ORANGE, YELLOW, DISCOVERY]):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = col
    ser.format.line.color.rgb = WHITE; ser.format.line.width = Pt(1)
for si, tier in enumerate([1, 2, 3, 4]):
    ser = ch.plots[0].series[si]
    for pi, proj in enumerate(("clean_4797", "dirty_4581")):
        pct = TIER[proj][tier][1]
        if pct < 8.0:
            continue
        lb = ser.points[pi].data_label
        lb.position = XL_LABEL_POSITION.CENTER
        lb.text_frame.text = f"{pct:.0f}%"
        pr = lb.text_frame.paragraphs[0]
        pr.font.size = Pt(11); pr.font.name = FB; pr.font.bold = True
        pr.font.color.rgb = WHITE if tier in (1, 4) else JET
ch.category_axis.tick_labels.font.size = Pt(11); ch.category_axis.tick_labels.font.name = FB
ch.value_axis.tick_labels.font.size = Pt(9); ch.value_axis.tick_labels.font.name = FB
ch.value_axis.has_major_gridlines = False
read_card(s, 8.42, 1.86, 4.31, 3.70, [
    "Blue on the right is Tier 4, meaning the record passed every check. Red on the left is Tier 1, meaning something in it is impossible.",
    "In the clean project almost 9 in 10 records pass. In the legacy project fewer than 1 in 20 do.",
    "That is why we treat the clean project as our main dataset and the legacy one as a rough second opinion.",
], stacked=True)
caution(s, "A low tier means we do not trust the record. It does not mean a bot filled it in. We have no way to tell who or what typed the answers.")
txt(s, M, 6.60, SW - 2 * M, 0.22, "Source: table_16_trust_tier_counts.csv", size=8, color=GREY, font=FB, space_after=0)
s.notes_slide.notes_text_frame.text = (
    "This is the clearest single picture of why we cannot just pool the two projects and enjoy the bigger sample.\n\n"
    "Every record in both projects goes through the same ten rules, and comes out with a tier from 1 to 4. The rules never look at the "
    "screening answer or anything else we care about, so the tiering cannot be tuned to make a result look nicer.\n\n"
    "In the clean project, eighty six percent land in Tier 4 and only two records out of 177 are Tier 1. In the legacy project it flips "
    "completely. Under four percent pass, and nearly five hundred records trip a hard impossibility.\n\n"
    "The sentence at the bottom is the one I most want people to repeat correctly. A Tier 1 record might be a rushed human, a form that "
    "misbehaved, or something automated. We genuinely cannot tell those apart, so we describe the evidence rather than accusing the respondent. "
    "If anyone asks me for a percentage of bots, the honest answer is that we do not have one and cannot get one from this data."
)

# ═════════════════════════════════════════════════════════════════════════════
# 8. R8 WORKED EXAMPLE
# ═════════════════════════════════════════════════════════════════════════════
s = blank(prs)
head2(s, "A simple check that catches impossible answers",
      "If the survey form could not have produced an answer, the record cannot be right. There is no threshold to argue about and no judgement call to make.")
# worked example, left
plate(s, M, 1.88, 6.55, 3.55, fill=COOLWHITE, radius=0.12)
txt(s, M + 0.28, 2.06, 6.0, 0.24, "A real example from record 4581_1069", size=9.5,
    color=ORANGE, font=FH, bold=True, caps=True, space_after=0, char_space=0.6)
# question 1
plate(s, M + 0.28, 2.40, 5.99, 0.70, fill=WHITE, radius=0.14)
txt(s, M + 0.50, 2.53, 4.30, 0.24, "How many children do you have?", size=12, color=JET,
    font=FB, space_after=0)
plate(s, M + 5.05, 2.52, 1.05, 0.46, fill=DISCOVERY, radius=0.22)
txt(s, M + 5.05, 2.605, 1.05, 0.30, "2", size=16, color=COOLWHITE, font=FH, bold=True,
    align=PP_ALIGN.CENTER, space_after=0)
# question 2
plate(s, M + 0.28, 3.22, 5.99, 0.92, fill=WHITE, radius=0.12)
txt(s, M + 0.50, 3.34, 5.55, 0.24, "Which age bands are your children in?", size=12,
    color=JET, font=FB, space_after=0)
for i, lab in enumerate(["Band 3", "Band 4", "Band 5"]):
    xx = M + 0.50 + i * 1.34
    plate(s, xx, 3.64, 1.20, 0.36, fill=ORANGE, radius=0.28)
    txt(s, xx, 3.715, 1.20, 0.24, lab, size=10.5, color=WHITE, font=FH, bold=True,
        align=PP_ALIGN.CENTER, space_after=0)
txt(s, M + 4.62, 3.685, 1.70, 0.24, "3 bands ticked", size=10.5, color=RED, font=FH,
    bold=True, space_after=0)
# verdict
plate(s, M + 0.28, 4.28, 5.99, 0.90, fill=RED, radius=0.12)
txt(s, M + 0.52, 4.40, 5.55, 0.70,
    "Three different age bands cannot hold two children.\nThe form does not allow it, so at least one answer here is not a real report.",
    size=11.5, color=WHITE, font=FB, space_after=0, line=1.05)
# numbers, right
NUMS = [("44", "records flagged by this check"),
        ("44", "of them in the legacy project"),
        ("0", "of them in the clean project"),
        ("0", "false alarms in 131 verified humans")]
for i, (big, lab) in enumerate(NUMS):
    yy = 1.88 + i * 0.92
    plate(s, 8.42, yy, 4.31, 0.80, fill=COOLBLUE if i % 2 == 0 else COOLWHITE, radius=0.13)
    txt(s, 8.68, yy + 0.13, 1.00, 0.52, big, size=27, color=DISCOVERY, font=FH, bold=True,
        space_after=0, char_space=-0.4)
    txt(s, 9.72, yy + 0.24, 2.85, 0.45, lab, size=10.5, color=JET, font=FB, space_after=0, line=1.0)
read_card(s, M, 5.62, 6.55, 1.16, [
    "Most fraud checks say a record looks unusual. This one says the record is impossible, which is much harder to argue with.",
    "It never flagged any of our 131 verified human caregivers, so tightening it costs no good data.",
], stacked=False, title="Why this one is different")
txt(s, M, 6.96, SW - 2 * M, 0.24,
    "Source: table_34_branching_logic_audit.csv (44 rows), table_15_rule_false_positive_rates.csv, table_15b_rule_counts_by_project.csv",
    size=8, color=GREY, font=FB, space_after=0)
s.notes_slide.notes_text_frame.text = (
    "This is the newest piece of the data quality work and I think it is the most defensible thing we have, so let me explain why with the "
    "example rather than the theory.\n\n"
    "A caregiver told us they have two children. Then, on the next question, they ticked three different age bands. Two children cannot occupy "
    "three age bands. The form was never designed to accept that, so at least one of those answers is not a truthful report.\n\n"
    "Compare that with our other checks. Most of them say something like this person finished faster than any verified human, which is "
    "suspicious but arguable. This one is not arguable. Either the form allows the combination or it does not, and it does not.\n\n"
    "Three numbers make the case. It caught forty four records. Every single one is in the legacy project and none are in the clean project, "
    "which is a small sanity check on the rule itself. And it flagged zero of our 131 verified human caregivers, so we can lean on it hard "
    "without losing any good data.\n\n"
    "One limitation I would volunteer before anyone asks. Forty four records is only about two and a half percent of the legacy project. "
    "This check is very accurate about what it catches, but it says nothing about how many other bad records are sitting there uncaught. "
    "It is not a contamination estimate."
)

# ═════════════════════════════════════════════════════════════════════════════
# 9. DOES DIRTY DATA CHANGE THE ANSWER
# ═════════════════════════════════════════════════════════════════════════════
simple_fig_slide(prs,
    "Does the messy data change the main finding?",
    "We re ran the whole analysis five times, using five different levels of strictness about which records to keep. The answer barely moved.",
    F("figure_20_tier_sensitivity.png"),
    ["Each row is one way of deciding which records to keep. The bottom row is the strictest, the top row uses only the legacy project.",
     "The dot is how big the gap between the two caregiver groups is. The line shows how sure we are about it.",
     "Every line sits well to the right of zero, so the finding holds whichever rule we pick. The gap ranges from about 42 to 52 points.",
     "The top row looks strongest but its line is very wide, because it rests on only 47 people. Treat it as consistent rather than as proof."],
    "figure_20_tier_sensitivity.png, table_20_tier_sensitivity.csv",
    "This is where the two halves of the talk meet. The fair worry is that the caregiver finding is an artefact of bad data, so we tested it directly.\n\n"
    "We picked five sensible ways of deciding which records deserve to be in the analysis, from only the very cleanest records right through to "
    "pooling in part of the legacy project. Then we re ran everything five times and looked at whether the gap between the two caregiver groups moved.\n\n"
    "It did not move much. It stays somewhere between forty two and fifty two points, and every one of those five lines stays clear of zero. "
    "If the messy records were manufacturing this result, tightening to the cleanest subset should have made it collapse. Instead it went up slightly.\n\n"
    "The structural reason is simple and worth saying. Inside the 131 caregivers we actually grouped, 117 already pass every check. There are only "
    "fourteen low trust records in there, so filtering does not have much to remove.\n\n"
    "I would not lean on the top row. It is the legacy only replication and it rests on forty seven people, which is why its line is so wide."
)

# ═════════════════════════════════════════════════════════════════════════════
# 10. WHAT WE CAN AND CANNOT SAY
# ═════════════════════════════════════════════════════════════════════════════
s = blank(prs)
head2(s, "What we can say, and what we will not say",
      "Everything on the left is supported by the figures we just walked through. Everything on the right is a sentence I would ask us not to write.")
plate(s, M, 1.96, 6.00, 4.10, fill=COOLBLUE, radius=0.14)
plate(s, 6.93, 1.96, 5.80, 4.10, fill=COOLWHITE, radius=0.14)
txt(s, M + 0.30, 2.16, 5.4, 0.30, "We can say this", size=15, color=DISCOVERY, font=FH,
    bold=True, space_after=0)
txt(s, 7.23, 2.16, 5.2, 0.30, "We will not say this", size=15, color=RED, font=FH,
    bold=True, space_after=0)
CAN = ["The two caregiver groups are about 42 points apart on whether they would definitely screen their baby.",
       "That gap survives once we account for autism in the home, and it holds under all five strictness settings we tried.",
       "44 records in the legacy project contain answers the survey form could not have produced.",
       "Almost 9 in 10 clean project records pass every trust check. Fewer than 1 in 20 legacy records do."]
CANT = ["That attitudes cause screening decisions. We measured both in the same sitting.",
        "That any share of records came from bots. We have no way to tell who typed them.",
        "That the two groups are the same on knowledge or demographics. We just could not detect a gap this small.",
        "That the 44 flagged records tell us how much of the legacy project is bad. They only tell us about those 44."]
for i, t in enumerate(CAN):
    yy = 2.62 + i * 0.85
    mk = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(M + 0.30), Inches(yy + 0.02), Inches(0.26), Inches(0.26))
    mk.fill.solid(); mk.fill.fore_color.rgb = DISCOVERY; mk.line.fill.background(); mk.shadow.inherit = False
    txt(s, M + 0.30, yy + 0.055, 0.26, 0.22, "Y", size=10, color=COOLWHITE, font=FH, bold=True,
        align=PP_ALIGN.CENTER, space_after=0)
    txt(s, M + 0.68, yy, 5.02, 0.78, t, size=11, color=JET, font=FB, space_after=0, line=1.03)
for i, t in enumerate(CANT):
    yy = 2.62 + i * 0.85
    mk = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.23), Inches(yy + 0.02), Inches(0.26), Inches(0.26))
    mk.fill.solid(); mk.fill.fore_color.rgb = RED; mk.line.fill.background(); mk.shadow.inherit = False
    txt(s, 7.23, yy + 0.055, 0.26, 0.22, "N", size=10, color=WHITE, font=FH, bold=True,
        align=PP_ALIGN.CENTER, space_after=0)
    txt(s, 7.61, yy, 4.82, 0.78, t, size=11, color=JET, font=FB, space_after=0, line=1.03)
plate(s, M, 6.26, SW - 2 * M, 0.62, fill=DISCOVERY, radius=0.28)
txt(s, M + 0.30, 6.42, SW - 2 * M - 0.60, 0.34,
    "Next steps: keep the clean project as our main dataset, report trust tiers instead of labels, and use the wording on the left when we write this up.",
    size=12, color=COOLWHITE, font=FH, bold=True, space_after=0)
source(s, "table_30_decision_summary.csv, plus every figure in this deck. Full detail is in ESD_Caregiver_Cluster_and_Trust_Screen.pptx")
s.notes_slide.notes_text_frame.text = (
    "I want to end here rather than on a result, because the wording is what actually leaves this room.\n\n"
    "The left column is what the figures support. Four sentences, and I would be comfortable defending any of them to a reviewer.\n\n"
    "The right column is the set of sentences that are tempting, natural to write, and not supported. The causal one is the easiest to slip "
    "into. The bot percentage is the one people will ask for directly, and I would rather we agree now that the answer is we do not have one. "
    "And the demographic equivalence claim is subtle, because a null result feels like evidence of sameness when it is really just evidence "
    "that our sample was too small to see anything.\n\n"
    "None of the right column weakens the finding. The gap is large, it survives adjustment, it survives every strictness setting, and the "
    "logic audit is clean. We just do not want to spend that credibility on claims the data cannot carry.\n\n"
    "If everyone is comfortable with the four sentences on the left, that is what goes into the manuscript draft, and the long deck has the "
    "traceability and the full checklist behind each one."
)

# ── hard constraint check: no em dashes, no en dashes in slide text ───────────
bad = []
for i, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        if sh.has_text_frame and ("—" in sh.text_frame.text or "–" in sh.text_frame.text):
            bad.append((i, sh.text_frame.text[:60]))
    if sl.has_notes_slide:
        nt = sl.notes_slide.notes_text_frame.text
        if "—" in nt or "–" in nt:
            bad.append((i, "NOTES: " + nt[:60]))
if bad:
    for b in bad:
        print("  DASH FOUND", b)
    raise SystemExit("Em or en dash present. Refusing to save.")
print("[dash check] no em dashes or en dashes in any slide text or speaker note")

prs.save(f"{OUTDIR}/ESD_Caregiver_Findings_Plain_Language.pptx")
print(f"[saved] plain-language deck, {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
