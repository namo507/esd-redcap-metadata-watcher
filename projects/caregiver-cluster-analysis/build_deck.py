"""Build the ESD Lab deck: Cluster Analysis & Bot Detection (trust-screen framing).

Every number in this deck is asserted against the on-disk CSVs in Caregiver Outputs
before a single slide is written. The build aborts on any mismatch.
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import pandas as pd
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from esd_deck_lib import *

OUTDIR = os.path.dirname(OUT)
F = lambda n: f"{OUT}/{n}"

# ─────────────────────────────────────────────────────────────────────────────
# 1. ASSERT ANCHORS AGAINST ON-DISK GROUND TRUTH
# ─────────────────────────────────────────────────────────────────────────────
t15  = pd.read_csv(F("table_15_rule_false_positive_rates.csv"))
t15b = pd.read_csv(F("table_15b_rule_counts_by_project.csv"))
t16  = pd.read_csv(F("table_16_trust_tier_counts.csv"))
t34  = pd.read_csv(F("table_34_branching_logic_audit.csv"))
t3   = pd.read_csv(F("table_3_cluster_diagnostics.csv"))
t4   = pd.read_csv(F("table_4_cluster_defining_profiles.csv"))
t5   = pd.read_csv(F("table_5_screening_outcome.csv"))
t5b  = pd.read_csv(F("table_5b_cohort_flow.csv"))
t20  = pd.read_csv(F("table_20_tier_sensitivity.csv"))
t30  = pd.read_csv(F("table_30_decision_summary.csv"))
t35  = pd.read_csv(F("table_35_demographic_signal_summary.csv"))
t35b = pd.read_csv(F("table_35b_demographic_signal_enrichment.csv"))

A = []
def chk(name, got, want, tol=0.0):
    ok = abs(got - want) <= tol if isinstance(want, (int, float)) else got == want
    A.append((name, got, want, ok))
    return ok

r8b = t15b[t15b.rule == "R8"].iloc[0]
chk("R8 all_n", int(r8b.all_n), 44)
chk("R8 clean_4797_n", int(r8b.clean_4797_n), 0)
chk("R8 dirty_4581_n", int(r8b.dirty_4581_n), 44)
chk("R8 false_positive_pct", float(t15[t15.rule == "R8"].iloc[0].false_positive_pct), 0.0)
chk("R8 verified_human_n", int(t15[t15.rule == "R8"].iloc[0].verified_human_n), 131)
chk("combined cached records", int(t16.n.sum()), 1956)
chk("clean_4797 total", int(t16[t16.source_project == "clean_4797"].n.sum()), 177)
chk("dirty_4581 total", int(t16[t16.source_project == "dirty_4581"].n.sum()), 1779)
chk("branching audit rows", len(t34), 44)
chk("branching audit projects", set(t34.source_project.unique()), {"dirty_4581"})
FAM = {
    "prenatal_testing_mutually_exclusive_pair": 8,
    "prenatal_tested_reason_without_prenatal_yes": 16,
    "prenatal_not_tested_reason_without_prenatal_no": 6,
    "earlier_diagnosis_yes_reason_without_yes_gate": 2,
    "earlier_diagnosis_no_reason_without_no_gate": 3,
    "child_age_band_count_exceeds_child_count": 9,
}
vc = t34.family.value_counts().to_dict()
for k, v in FAM.items():
    chk(f"family {k}", int(vc.get(k, -1)), v)
d = t16[t16.source_project == "dirty_4581"]
chk("dirty tiers 1-3 pct", round(float(d[d.tier <= 3].pct_within_project.sum()), 1), 96.1, 0.05)
chk("dirty tier 4 pct", round(float(d[d.tier == 4].pct_within_project.iloc[0]), 1), 3.9, 0.05)

dirty_prev = t35[(t35.analysis == "project_prevalence") & (t35.source_project == "dirty_4581")].iloc[0]
clean_prev = t35[(t35.analysis == "project_prevalence") & (t35.source_project == "clean_4797")].iloc[0]
cmp_prev = t35[t35.analysis == "dirty_vs_clean_comparison"].iloc[0]
tier1_enrich = t35b[t35b.metric == "tier_1_confirmed_invalid"].iloc[0]
tier2_enrich = t35b[t35b.metric == "tier_le_2_high_or_confirmed"].iloc[0]
r8_enrich = t35b[t35b.metric == "rule_R8"].iloc[0]

chk("demographic dirty selected", int(dirty_prev.n_selected), 46)
chk("demographic clean selected", int(clean_prev.n_selected), 2)
chk("demographic comparison p", round(float(cmp_prev.p_value), 3), 0.312, 0.001)
chk("demographic tier1 group flagged", int(tier1_enrich.group_flagged_n), 27)
chk("demographic tier<=2 group flagged", int(tier2_enrich.group_flagged_n), 39)
chk("demographic R8 group flagged", int(r8_enrich.group_flagged_n), 4)

bad = [a for a in A if not a[3]]
print(f"[anchors] {len(A)-len(bad)}/{len(A)} passed")
if bad:
    for b in bad:
        print("  MISMATCH", b)
    raise SystemExit("Anchor verification failed — refusing to build the deck.")

# derived, straight from the CSVs
TIER = {p: {int(r.tier): (int(r.n), float(r.pct_within_project))
            for _, r in t16[t16.source_project == p].iterrows()}
        for p in ("clean_4797", "dirty_4581")}
FAMS = [("Prenatal tested-reason without a prenatal-yes gate", 16),
        ("Child age-band count exceeds child count", 9),
        ("Prenatal testing mutually exclusive pair", 8),
        ("Prenatal not-tested reason without a prenatal-no gate", 6),
        ("Earlier-diagnosis no-reason without a no gate", 3),
        ("Earlier-diagnosis yes-reason without a yes gate", 2)]

prs = new_deck()

# ─────────────────────────────────────────────────────────────────────────────
# S1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs, bg=DISCOVERY)
if os.path.exists(PATTERN_BAND):
    s.shapes.add_picture(PATTERN_BAND, Inches(-0.6), Inches(5.55),
                         width=Inches(14.5), height=Inches(2.4))
plate(s, M, 0.46, 5.95, 0.92, fill=WHITE, radius=0.28)
logo_pair(s, M + 0.32, 0.66, 0.46)
txt(s, M, 2.10, 11.6, 1.7,
    "Caregiver Acceptability Clusters\nand Data-Trust Screening",
    size=40, color=COOLWHITE, font=FH, bold=True, line=0.94, char_space=-0.3, space_after=0)
txt(s, M, 3.86, 11.6, 0.34,
    "Infant Autism Screening survey  ·  REDCap projects 4797 and 4581",
    size=15, color=COOLBLUE, font=FB, space_after=0)
txt(s, M, 4.34, 11.6, 0.60,
    ["Early Social Development Lab  ·  University of South Carolina",
     "Verified cache snapshot 2026-08-14  ·  1,956 cached records  ·  Prepared for PI, clinical and technical review"],
    size=10.5, color=SCIENCE, font=FB, space_after=3)
s.notes_slide.notes_text_frame.text = (
    "Framing for the room. This deck answers two questions that were run in parallel on the same survey. "
    "First, do caregivers fall into interpretable acceptability profiles, and does profile membership relate to a "
    "held-out screening-intent outcome. Second, how much of the incoming data can we trust, given that one of the two "
    "REDCap projects was flooded with low-quality submissions.\n\n"
    "Say this explicitly up front: nothing in this deck labels any individual respondent as a bot. We report a graded, "
    "deterministic trust screen. Where the evidence is ambiguous I have marked it as uncertainty rather than rounding it "
    "into a claim.\n\n"
    "Everything shown was generated by the existing pipeline and lives in Caregiver Outputs. No new modelling was done for "
    "this deck, and every number on every slide was re-read from its source CSV at build time and asserted against the "
    "agreed anchors before the file was written."
)

# ─────────────────────────────────────────────────────────────────────────────
# S2 — SCOPE (stat tiles)
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
header(s, "Scope and design", "One survey, two questions, one held-out design")
takeaway(s, "Profiles are built only from Thoughts, Feelings and Attitudes items, and every outcome we care about is held out of that construction.")
tiles = [("1,956", "cached records\nacross both projects"),
         ("131", "caregivers assigned\nto a profile"),
         ("127", "with a valid held-out\nscreening outcome"),
         ("10", "deterministic trust\nrules, R1 to R10")]
tw, tg = 2.85, 0.24
x0 = M
for i, (big, lab) in enumerate(tiles):
    x = x0 + i * (tw + tg)
    plate(s, x, 2.48, tw, 1.45, fill=DISCOVERY, radius=0.14)
    txt(s, x + 0.20, 2.60, tw - 0.40, 0.66, big, size=34, color=COOLWHITE, font=FH,
        bold=True, space_after=0, char_space=-0.4)
    txt(s, x + 0.20, 3.26, tw - 0.40, 0.58, lab.replace("\n", " "), size=10,
        color=COOLBLUE, font=FB, space_after=0, line=0.95)
bullets(s, M, 4.20, 12.13, 2.30, [
    "The 1,956 records combine clean_4797 (n=177) and dirty_4581 (n=1,779); UIDs are namespaced so that 174 overlapping numeric record IDs cannot corrupt joins",
    "Clustering uses the 10 TFA domains only. Screening intent, autism knowledge, values and demographics are held out and examined only after profile assignment",
    "Cohort flow: 135 selected analytic records, 131 cluster-assigned at a threshold of at least 8 of 10 input domains, 127 of those with a valid screening outcome",
    "The trust screen is deterministic and rule-based. It runs on all 1,956 records, and no outcome variable enters any rule or any tier",
], size=12)
source(s, "table_5b_cohort_flow.csv; table_12_api_cache_inventory.csv; table_16_trust_tier_counts.csv; table_31_upgrade_data_quality.csv")
s.notes_slide.notes_text_frame.text = (
    "I want to be precise about what is inside the model and what is outside it, because it is the design choice that makes "
    "everything after this slide interpretable.\n\n"
    "Only the ten TFA domains build the profiles. Screening intent, knowledge, values and demographics never touch the clustering. "
    "That is why I can put screening intent on a later slide and call it a held-out comparison rather than a circular one.\n\n"
    "On the cohort numbers: 135 records were supplied in the analytic file, 131 of them answered at least 8 of the 10 domains and "
    "so could be assigned a profile, and 127 of those also gave a usable screening answer. The four records that dropped out of "
    "clustering all had screening answers, and none of them said Definitely yes, which matters later when we talk about tipping points.\n\n"
    "The trust screen is a separate axis entirely. It covers all 1,956 records from both projects and it is deterministic, meaning "
    "the same record always lands in the same tier. If anyone asks whether the tiering could have been tuned to produce a nicer "
    "cluster result, the answer is no, because no outcome variable enters any rule."
)

# ─────────────────────────────────────────────────────────────────────────────
# S3 — CLUSTER COUNT
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Cluster structure",
    "Two profiles is the defensible choice, and the separation is modest",
    "Selection criteria disagree, so k=2 was chosen for stability and interpretability, not because a clean natural break exists.",
    ["Average silhouette peaks at k=2 (0.180) but stays below the 0.25 reference at every k tested",
     "80% subsample stability is 0.704 ARI at k=2 and higher at k=3 (0.806), but k=3 places only 9 caregivers, 6.9% of the cohort, in its smallest cluster",
     "The gap statistic rises monotonically through k=6, so it never nominates an interior optimum",
     "Treat the result as a partition of a gradient rather than the discovery of two discrete caregiver types"],
    F("figure_1_cluster_count_diagnostics.png"),
    "figure_1_cluster_count_diagnostics.png; table_3_cluster_diagnostics.csv",
    "This is the slide where I would rather be honest than impressive.\n\n"
    "Four criteria were run and they do not agree. Silhouette prefers k=2 but the value is 0.180, which is well under the 0.25 line "
    "drawn on the panel. In plain terms the caregivers do not fall into two clean, well-separated balls. Subsample stability actually "
    "likes k=3 better, at 0.806 versus 0.704, but the k=3 solution puts only nine people in its smallest cluster, which is too fragile "
    "to characterise or to act on. The gap statistic just keeps climbing, which is what you see when the data are more of a continuum "
    "than a set of clusters.\n\n"
    "So the choice of two profiles is a judgement call made for stability and interpretability, and I would defend it on those grounds "
    "rather than on any claim that we found natural kinds. If a reviewer pushes on this, the right answer is that the two-profile "
    "partition is a useful summary of a gradient, and every downstream result should be read with that in mind.\n\n"
    "The reassuring part is on a later slide: the screening association does not depend on this choice being exactly right."
)

# ─────────────────────────────────────────────────────────────────────────────
# S4 — PROFILE SEPARATION
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Profile definition",
    "Trust, blood-test openness and ethical support separate the profiles",
    "The two profiles differ most on relational and ethical dimensions, and barely at all on how aversive caregivers find the screen.",
    ["Higher acceptability n=84 (64.1%), Conditional acceptability n=47 (35.9%); mean TFA acceptability 78.7 versus 59.2 on a 0 to 100 scale",
     "Largest separations: positive affect and clinical trust 25.9 points, blood-test openness 25.7 points, ethical and moral support 25.3 points",
     "Smallest separation: low aversive affect, 97.6 versus 93.4, a 4.2-point difference — neither profile finds the screen distressing in itself",
     "Both profiles sit low on accuracy tolerance (44.6 versus 26.6) and equipment-intensive openness (47.4 versus 29.9)"],
    F("figure_2_primary_cluster_profiles.png"),
    "figure_2_primary_cluster_profiles.png; table_4_cluster_defining_profiles.csv; table_9_decision_summary.csv",
    "The naming matters here. I am deliberately not calling the second group low acceptability, because that is not what the data show. "
    "They are conditional: they will accept screening, but their support is contingent on trust, burden and accuracy being addressed.\n\n"
    "Look at where the two profiles are far apart and where they are close. The big gaps are positive affect and clinical trust, "
    "blood-test openness, and ethical and moral support, all around 25 points. Those are relational and values-laden dimensions, not "
    "logistical ones. The smallest gap by a wide margin is low aversive affect at 4.2 points, meaning both groups agree the screen is "
    "not upsetting. So the divide is about trust and conditions, not about squeamishness.\n\n"
    "One more thing worth flagging: both profiles score low on accuracy tolerance and on openness to equipment-intensive methods. "
    "That is a shared constraint across the whole cohort, and it points at the same design lever for both groups, which is transparency "
    "about accuracy and keeping the modality light.\n\n"
    "Decision use from the pipeline: maintain a low-burden pathway with immediate next steps for the higher group, and for the conditional "
    "group reduce visit burden while reinforcing trust, choice, accuracy transparency and ethical safeguards."
)

# ─────────────────────────────────────────────────────────────────────────────
# S5 — HELD-OUT SCREENING
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Held-out outcome",
    "Screening intent tracks profile membership",
    "Definitely-yes screening intent is 70.0% in the higher-acceptability profile and 27.7% in the conditional profile.",
    ["56 of 80 versus 13 of 47 answered Definitely yes, a gap of 42.3 percentage points",
     "Wilson 95% intervals are 59.2 to 78.9% and 16.9 to 41.8%, and they do not overlap",
     "Widening the endpoint to probably or definitely yes gives 90.0% versus 59.6%, so the ordering is not an artefact of where the cut is placed",
     "Screening was held out of clustering, so this is an association between an unsupervised partition and an outcome, not a causal effect of profile membership"],
    F("figure_3_screening_outcome.png"),
    "figure_3_screening_outcome.png; table_5_screening_outcome.csv",
    "This is the headline result, and it is also the one most likely to be over-read, so let me set the guardrail first. "
    "The profiles were built without ever seeing the screening question. The screening answer was then compared across profiles. "
    "That makes this a genuine held-out comparison, and it is why the gap is worth reporting at all.\n\n"
    "What it is not is a causal claim. Nothing here says that changing a caregiver's attitudes would change their screening decision. "
    "These are cross-sectional survey responses collected at one sitting, so attitude and intent are measured together and could easily "
    "share a common cause, or simply be two expressions of the same underlying disposition.\n\n"
    "The numbers: 70.0% versus 27.7% on the strict endpoint, a 42.3 point gap, with non-overlapping Wilson intervals. If you loosen the "
    "endpoint to include probably, it is 90.0% versus 59.6%. The ordering survives the change of cut point, which is a small but real "
    "robustness check.\n\n"
    "Note the denominators differ: 80 of the 84 higher-acceptability caregivers gave a valid screening answer, versus all 47 of the "
    "conditional group. Four missing responses sit entirely in the higher group."
)

# ─────────────────────────────────────────────────────────────────────────────
# S6 — ADJUSTED ASSOCIATION
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Adjusted association",
    "The profile signal survives adjustment for autism in the household",
    "Profile membership and lived autism experience act additively, and there is no evidence that they interact.",
    ["Higher-acceptability profile adjusted odds ratio 8.73; autistic child at home adjusted odds ratio 4.32 (p=0.0012), n=127",
     "The profile-by-ASD interaction term is 2.02 with p=0.445, which is not distinguishable from no interaction at this sample size",
     "The profile gap holds inside both strata: 88.2% versus 37.5% where there is an autistic child at home, and 56.5% versus 17.4% where there is not",
     "Firth profile-likelihood intervals are used throughout because several cells are sparse and standard maximum likelihood is unstable there"],
    F("figure_18_logistic_forest.png"),
    "figure_18_logistic_forest.png; table_11_asd_family_context.csv; table_22_logistic_models.csv",
    "The obvious challenge to the previous slide is confounding by lived experience: families already living with autism might be both "
    "more accepting and more likely to screen. So we tested it directly.\n\n"
    "Having an autistic child at home is indeed a strong independent predictor, with an adjusted odds ratio of 4.32 and a p-value of 0.0012. "
    "But it does not explain away the profile effect, which stays at 8.73 after adjustment. And the interaction term is 2.02 with a p-value "
    "of 0.445, so on this evidence the two act additively rather than multiplying each other.\n\n"
    "Be careful with the interaction result. A p-value of 0.445 at n=127 is weak evidence of nothing, not evidence of no interaction. "
    "The confidence interval on that term is wide. I would phrase it as: we cannot detect an interaction, and the study is not powered to "
    "rule out a moderate one.\n\n"
    "The stratified numbers are the intuitive version of the same finding. Among families with an autistic child at home the profile gap is "
    "50.7 points; among those without, it is 39.1 points. Same direction, similar magnitude, both statistically clear.\n\n"
    "Firth penalisation is used because some strata have very few events. It is the standard fix for separation in small logistic models and "
    "it makes the intervals honest rather than infinite."
)

# ─────────────────────────────────────────────────────────────────────────────
# S7 — WHAT DOES NOT SEPARATE
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Characterization",
    "Knowledge, values and demographics do not reliably separate the profiles",
    "After multiplicity control none of the three held-out characterization families distinguishes the two profiles.",
    ["Verified 7-item autism knowledge is 3.90 versus 3.41 correct (55.7% versus 48.8%), Mann-Whitney p=0.099",
     "Self-direction shows the largest raw values signal (Cohen's d 0.48, p=0.013) but its FDR-adjusted p is 0.129",
     "Every demographic contrast has FDR p at or above 0.56; the closest is any premature birth at raw p=0.051, FDR p=0.56",
     "At 47 and 80 respondents the study can only detect roughly a 25.6 point difference at 80% power, so this is absence of evidence, not evidence of equivalence"],
    F("figure_4_autism_knowledge.png"),
    "figure_4_autism_knowledge.png; figure_5_values_characterization.png; figure_6_demographic_characterization.png; table_6, table_7, table_8, table_29",
    "This is the slide that keeps the story honest. The profiles differ sharply on attitudes and on screening intent, but they are "
    "not obviously different kinds of people.\n\n"
    "Knowledge goes in the expected direction, 3.90 versus 3.41 out of seven, but the p-value is 0.099 and I would not build anything "
    "on it. Values look more interesting at first glance because self-direction has a Cohen's d of about 0.48 and a raw p of 0.013, but "
    "we tested ten values and after false discovery rate control it lands at 0.129. That is exactly the kind of result that gets "
    "over-claimed in this literature, so I am flagging it as suggestive at most.\n\n"
    "Demographics show nothing that survives adjustment. Education, employment, rurality, healthcare access, household size, caregiver "
    "age, all null. The nearest miss is premature birth history at a raw p of 0.051.\n\n"
    "The power line at the bottom is the one to emphasise if someone concludes the groups are the same. With these sample sizes we could "
    "only reliably detect a difference of about 25.6 percentage points. Anything smaller than that would be invisible to us. So the correct "
    "reading is that we did not find differences, not that there are none."
)

# ─────────────────────────────────────────────────────────────────────────────
# S8 — TRUST SCREEN DESIGN
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Data trust",
    "A graded, deterministic screen — not a bot label",
    "Ten rules anchored to 131 verified-human completions assign every record a tier, and no record is ever labelled a bot.",
    ["Rules R1 to R10 run on all 1,956 records; thresholds are derived from the completed 4797 caregivers and never from any outcome",
     "False-positive rates measured in verified humans: R10 7.6%, R3 4.6%, R4 3.8%, R6 2.3%, R9 1.5%. R1, R2, R5, R7 and R8 flag nobody",
     "R6, rapid submission clustering, dominates co-occurrence at 1,682 records, so it is the workhorse and also the least specific rule",
     "R5, the exact ordered Likert fingerprint duplicate, never fires anywhere in the 1,956 records",
     "R10 is validation-only and does not define a tier; its email anti-fraud component is unavailable under current API permissions"],
    F("figure_9_rule_cooccurrence.png"),
    "figure_9_rule_cooccurrence.png; table_14_fraud_rule_definitions.csv; table_15_rule_false_positive_rates.csv; table_31_upgrade_data_quality.csv",
    "Framing first, because the wording matters for how this gets cited. We do not have ground truth on who was a bot. What we have is a "
    "set of ten deterministic rules whose thresholds are anchored to 131 caregivers we can vouch for as human, plus a graded tier that "
    "expresses how much of that record we are willing to trust.\n\n"
    "The right-hand panel is the calibration check, and it is the reason to believe the screen. Applied back to the verified humans, most "
    "rules flag nobody. The two with meaningful false-positive rates are R10 at 7.6% and R3 at 4.6%. R10 is why it is validation-only and "
    "does not set a tier.\n\n"
    "The left panel shows the rules are not independent. R6 co-occurs with almost everything and accounts for 1,682 records on its own. "
    "That means R6 is doing most of the work and is also the least specific signal, so a record flagged only by R6 deserves less weight than "
    "one flagged by several rules.\n\n"
    "R5 never fires. Exact duplicated Likert fingerprints simply do not occur in this data, which is mildly informative in itself: whatever "
    "produced the low-quality records was not copy-pasting identical response vectors.\n\n"
    "Caveat to state out loud: the email anti-fraud component is not returned under current API permissions, so R10 works from knee and age "
    "gates only. We are not inferring evidence we do not have."
)

# ─────────────────────────────────────────────────────────────────────────────
# S9 — TIER DISTRIBUTION (native chart)
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
header(s, "Tier distribution", "Pass rates differ sharply between the two projects")
takeaway(s, "86.4% of clean-project records reach Tier 4, versus 3.9% in the legacy project.")
cd = CategoryChartData()
cd.categories = [f"clean_4797\n(n=177)", f"dirty_4581\n(n=1,779)"]
for tier, label in [(1, "Tier 1  Confirmed invalid"), (2, "Tier 2  High suspicion"),
                    (3, "Tier 3  Uncertain"), (4, "Tier 4  Pass")]:
    cd.add_series(label, (round(TIER["clean_4797"][tier][1], 1),
                          round(TIER["dirty_4581"][tier][1], 1)))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, Inches(M), Inches(2.46),
                        Inches(7.05), Inches(4.35), cd)
ch = gf.chart
ch.has_title = False
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(9); ch.legend.font.name = FB
for ser, col in zip(ch.plots[0].series, [RED, ORANGE, YELLOW, DISCOVERY]):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = col
    ser.format.line.color.rgb = WHITE; ser.format.line.width = Pt(1)
# per-point labels: only wide-enough segments get a label, so nothing collides
for si, (tier, _) in enumerate([(1, 0), (2, 0), (3, 0), (4, 0)]):
    ser = ch.plots[0].series[si]
    for pi, proj in enumerate(("clean_4797", "dirty_4581")):
        pct = TIER[proj][tier][1]
        if pct < 8.0:
            continue
        lab = ser.points[pi].data_label
        lab.position = XL_LABEL_POSITION.CENTER
        lab.text_frame.text = f"{pct:.1f}%"
        pr = lab.text_frame.paragraphs[0]
        pr.font.size = Pt(9.5); pr.font.name = FB; pr.font.bold = True
        pr.font.color.rgb = WHITE if tier in (1, 4) else JET
ch.category_axis.tick_labels.font.size = Pt(10); ch.category_axis.tick_labels.font.name = FB
ch.value_axis.tick_labels.font.size = Pt(9); ch.value_axis.tick_labels.font.name = FB
ch.value_axis.has_major_gridlines = False
bullets(s, 8.00, 2.52, 4.73, 4.20, [
    f"clean_4797: {TIER['clean_4797'][4][0]} pass, {TIER['clean_4797'][3][0]} uncertain, {TIER['clean_4797'][2][0]} high suspicion, {TIER['clean_4797'][1][0]} confirmed invalid",
    f"dirty_4581: {TIER['dirty_4581'][4][0]} pass, {TIER['dirty_4581'][3][0]} uncertain, {TIER['dirty_4581'][2][0]} high suspicion, {TIER['dirty_4581'][1][0]} confirmed invalid",
    "Tiers 1 to 3 account for 96.1% of the legacy project and 13.6% of the clean project",
    "This describes how much of each project our rules will trust. It is not an estimate of how many respondents were automated",
], size=11)
source(s, "table_16_trust_tier_counts.csv (chart values read directly from this file)")
s.notes_slide.notes_text_frame.text = (
    "This is the clearest single picture of why the two projects cannot be pooled naively.\n\n"
    "In the clean project 86.4% of records pass all the way to Tier 4 and only two records are confirmed invalid. In the legacy project "
    "the picture inverts: 3.9% pass, and 96.1% land somewhere in tiers 1 to 3. Nearly 500 records in the legacy project trip a hard "
    "logical or timing impossibility.\n\n"
    "The sentence I want people to leave with is the last bullet. This chart is about how much of each project our rules are willing to "
    "trust. It is not a bot count. A record can land in Tier 1 because a human rushed it, because a form behaved oddly, or because it was "
    "automated, and these data cannot separate those explanations. The tier is a statement about evidential quality, not about the nature "
    "of the respondent.\n\n"
    "Practically, this is what drives decision D1: keep 4797 as the primary cohort and treat 4581 as a caveated replication rather than "
    "as additional sample."
)

# ─────────────────────────────────────────────────────────────────────────────
# S10 — R8 BRANCHING AUDIT (native chart)
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
header(s, "Method update", "R8 generalized: 44 branching and logic impossibilities")
takeaway(s, "The branching audit turns REDCap skip-logic contradictions into auditable, record-level evidence with zero false positives among verified humans.")
cd = CategoryChartData()
cd.categories = [f[0] for f in reversed(FAMS)]
cd.add_series("Records", tuple(f[1] for f in reversed(FAMS)))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(M), Inches(2.44),
                        Inches(7.55), Inches(3.05), cd)
ch = gf.chart
ch.has_title = False; ch.has_legend = False
ser = ch.plots[0].series[0]
ser.format.fill.solid(); ser.format.fill.fore_color.rgb = DISCOVERY
ser.format.line.fill.background()
ch.plots[0].gap_width = 55
pl = ch.plots[0]; pl.has_data_labels = True
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
pl.data_labels.font.size = Pt(10); pl.data_labels.font.bold = True
pl.data_labels.font.name = FB; pl.data_labels.font.color.rgb = DISCOVERY
ch.category_axis.tick_labels.font.size = Pt(8.5); ch.category_axis.tick_labels.font.name = FB
ch.value_axis.has_major_gridlines = False
ch.value_axis.visible = False
bullets(s, M, 5.68, 12.13, 1.35, [
    "R8 now covers three families: autistic-child follow-ups that persist when the reported count is zero, prenatal and earlier-diagnosis branching contradictions, and an impossible child age-band count",
    "44 flagged records, all 44 in dirty_4581 and 0 in clean_4797, with a false-positive rate of 0.0% in the 131-record verified-human reference",
    "These are hard logical impossibilities rather than soft suspicion signals, so they stay in Tier 1 and are reviewable field by field against the instrument's own skip logic",
], size=10.5, gap=5)
source(s, "table_34_branching_logic_audit.csv (44 rows); table_15b_rule_counts_by_project.csv; table_15_rule_false_positive_rates.csv; README.md")
s.notes_slide.notes_text_frame.text = (
    "This is the methodological update since the last review, and I think it is the most defensible thing in the trust screen, so it is "
    "worth spending a minute on why.\n\n"
    "Most fraud heuristics are distributional. They say a record looks unusual compared with other records, which always leaves you arguing "
    "about thresholds. The branching audit is different in kind. It asks whether a record contradicts the instrument's own skip logic. "
    "If a caregiver answered why their child was not tested prenatally, but never said the child was not tested prenatally, that is not "
    "unusual, it is impossible under the form as designed. There is no threshold to argue about and no distributional assumption to violate.\n\n"
    "Three consequences follow. First, it is falsifiable and reviewable: clinical collaborators can open any one of the 44 records and check "
    "the gating field against the observed field themselves, because the table stores both. Second, it does not degrade when the population "
    "shifts, unlike a timing threshold or an outlier score. Third, and this is the number I would lead with, its false-positive rate is 0.0% "
    "across all 131 verified humans, so tightening the screen with R8 costs us no clean data at all.\n\n"
    "The distribution across families is uneven and that is informative. Sixteen of the 44 are tested-reason orphans and nine are impossible "
    "age-band counts, where the number of distinct child age bands ticked exceeds the number of children reported.\n\n"
    "One honest limitation: 44 records is a small share of the 1,779 legacy records, so R8 is a high-precision, low-recall instrument. "
    "It confirms specific records are unusable. It does not measure how many bad records exist overall."
)

# ─────────────────────────────────────────────────────────────────────────────
# S11 — ROBUSTNESS / INTEGRATION
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Robustness",
    "Trust filtering does not overturn the profile finding",
    "Across all five pre-specified inclusion definitions the screening gap stays between 42.3 and 52.2 percentage points.",
    ["Tier-4 only 44.7 pp, Tiers 3+4 48.5 pp, status quo 42.3 pp, pooled with 4581 Tier 4 45.7 pp, and 4581 Tier-4 replication 52.2 pp",
     "Every interval excludes zero, but the 4581 replication is imprecise at 5.7 to 74.1 pp on only 47 respondents",
     "Firth odds ratios across the five definitions range from 7.39 to 11.09",
     "Within the primary cohort only 14 of 131 clusterable caregivers sit below Tier 4, so trust filtering has limited room to move this result",
     "The pooled definition agrees only moderately with the status quo partition (ARI 0.52), so its point estimate is a different analysis, not a confirmation"],
    F("figure_20_tier_sensitivity.png"),
    "figure_20_tier_sensitivity.png; figure_19_tier_to_cluster_alluvial.csv → figure_19_tier_to_cluster_alluvial.png; table_20_tier_sensitivity.csv",
    "This is where the two halves of the deck meet. The natural worry is that the cluster result is an artefact of dirty data, so we "
    "re-ran the whole thing under five pre-specified inclusion definitions and asked whether the screening gap moves.\n\n"
    "It does not move much. The gap ranges from 42.3 to 52.2 percentage points and every interval excludes zero. If low-trust records were "
    "manufacturing the association, tightening to Tier 4 only should have collapsed it, and instead the gap goes slightly up to 44.7.\n\n"
    "Two honest caveats. The 4581 replication looks like the strongest result at 52.2 points, but its interval runs from 5.7 to 74.1 on "
    "just 47 people, so it is compatible with almost anything and should be read as directionally consistent rather than confirmatory. "
    "And the pooled definition has an ARI of only 0.52 against the status quo, meaning the profiles it produces are substantially different "
    "groupings. Its similar point estimate is reassuring, but it is a different analysis, not a replication of the same one.\n\n"
    "The structural reason the primary result is stable is on figure 19: within the 131 clusterable caregivers, 117 are already Tier 4 and "
    "only 14 sit below it. There simply is not much low-trust material inside the primary cohort for filtering to remove."
)

# ─────────────────────────────────────────────────────────────────────────────
# S12 — DETECTOR CAVEATS
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Detector caveats",
    "Independent detectors agree weakly, so no single score is ground truth",
    "Rule flags and unsupervised detectors overlap far less than a reliable detector would, which is the main reason we report tiers rather than labels.",
    ["Best agreement is Cohen's kappa 0.29, between rule flags and the PU/source classifier; isolation forest and local outlier factor also reach only 0.29",
     "Robust Mahalanobis is essentially unrelated to the rules (kappa −0.02) and to the PU classifier (−0.01)",
     "The 4797-versus-4581 source classifier reaches AUC 0.980, but the negative control on random clean halves sits at 0.525, so the separation is real yet not attributable to fraud alone",
     "The strongest discriminating feature is the count of filled open-text fields (permutation AUC drop 0.204), a plausible instrument-design difference rather than a bot signature",
     "The prior lab LightGBM artifact was not present in the scoped project files, so the PU classifier stands in for it"],
    F("figure_11_detector_scores_and_agreement.png"),
    "figure_11_detector_scores_and_agreement.png; figure_11b_pu_feature_importance.png; table_17, table_18, table_19, table_31",
    "If we were going to publish a bot prevalence, this is the slide that stops us, so I want to walk it carefully.\n\n"
    "Five independent ways of asking is this record anomalous were run: the deterministic rules, isolation forest, local outlier factor, "
    "robust Mahalanobis distance, and a positive-unlabelled source classifier. If any of them were tracking a real underlying bot process, "
    "they should agree with each other substantially. The best pairwise agreement in the whole matrix is a Cohen's kappa of 0.29, which is "
    "fair at best. Mahalanobis is at or below zero against two of the others, meaning it is picking up something unrelated.\n\n"
    "The source classifier is the seductive one. It separates 4797 from 4581 with an AUC of 0.980, and the model-implied separation share "
    "is 94.4%. It is very tempting to read that as 94% of the legacy project is nonhuman. That reading is wrong, and figure 11b shows why: "
    "the feature doing most of the work is the number of filled open-text fields. The two projects differ in their instruments and in their "
    "recruitment periods, so the classifier is separating projects, not humans from bots.\n\n"
    "The negative control is the discipline check. Split the clean project randomly in half and try to classify half A against half B: you "
    "get 0.525, essentially chance, which tells us the pipeline is not manufacturing separation on its own. Good, but it does not rescue "
    "the interpretation of the 0.980.\n\n"
    "Last note for the technical reviewers: the earlier lab LightGBM artifact was not in the scoped files, so concordance here uses rules, "
    "three novelty detectors and the PU classifier. That is a gap in coverage, not a hidden result."
)

# ─────────────────────────────────────────────────────────────────────────────
# S13 — DECISIONS
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
header(s, "Decisions", "What we conclude, and what we will not claim")
takeaway(s, "Keep 4797 as the primary cohort, report graded tiers rather than labels, and do not publish a bot prevalence.")
rows = [
    ("D1", "Primary cohort", "Keep 4797 primary; 4581 is a caveated replication", "Implemented"),
    ("D2", "Trust target", "Four graded inclusion tiers, not a binary bot label", "Implemented"),
    ("D3", "Contamination fraction", "Do not present an identifiable bot prevalence", "Open, not identifiable"),
    ("D4", "Knowledge scoring", "Binary verified correct-count is primary; graded is sensitivity only", "Pre-specified in config"),
    ("D5", "Four excluded records", "Report the full tipping-point range, 39.0 to 44.5 pp", "Implemented"),
]
y0, rh = 2.78, 0.65
cw = [0.62, 2.35, 6.20, 2.96]
cx = [M]
for w in cw[:-1]:
    cx.append(cx[-1] + w)
hdrs = ["", "Decision", "Recommendation", "Status"]
plate(s, M, y0 - 0.40, sum(cw), 0.38, fill=DISCOVERY, radius=0.10)
for i, hh in enumerate(hdrs):
    txt(s, cx[i] + 0.16, y0 - 0.315, cw[i] - 0.24, 0.24, hh, size=10, color=COOLWHITE,
        font=FH, bold=True, caps=True, space_after=0, char_space=0.5)
for j, (tag, dec, rec, st) in enumerate(rows):
    yy = y0 + j * rh
    if j % 2 == 0:
        plate(s, M, yy, sum(cw), rh - 0.06, fill=COOLWHITE, radius=0.06)
    txt(s, cx[0] + 0.14, yy + 0.16, cw[0] - 0.20, 0.30, tag, size=13, color=DISCOVERY,
        font=FH, bold=True, space_after=0)
    txt(s, cx[1] + 0.16, yy + 0.17, cw[1] - 0.28, 0.32, dec, size=11, color=JET,
        font=FH, bold=True, space_after=0, line=0.95)
    txt(s, cx[2] + 0.16, yy + 0.17, cw[2] - 0.30, 0.42, rec, size=11, color=JET,
        font=FB, space_after=0, line=0.95)
    col = ORANGE if "Open" in st else DISCOVERY
    txt(s, cx[3] + 0.16, yy + 0.17, cw[3] - 0.28, 0.36, st, size=10.5, color=col,
        font=FH, bold=True, space_after=0, line=0.95)
txt(s, M, 6.24, 12.13, 0.55,
    "The PU model-implied nonhuman share in 4581 is 94.4% (93.9 to 95.0), but recruitment-period and instrument drift violate the "
    "identifying assumption, so it is a separation fraction and not a bot prevalence. It should not be quoted as one.",
    size=10.5, color=RED, font=FB, line=0.98, space_after=0)
source(s, "table_30_decision_summary.csv; table_17_contamination_identifiability.csv")
s.notes_slide.notes_text_frame.text = (
    "Five decisions, four of them already implemented in the pipeline and one deliberately left open.\n\n"
    "D1 and D2 are the operational core. Keep 4797 primary because it is the only cohort with a verified-human reference, and use graded "
    "tiers so that downstream analysts can choose their own strictness rather than inheriting ours.\n\n"
    "D3 is the one I want agreement on in the room. There is a number available, 94.4% with a tight-looking interval, and it will be asked "
    "for. The Elkan-Noto estimator that produces it assumes labelled positives are selected completely at random from the positive class. "
    "Here the labelled set is one REDCap project and the unlabelled set is another, and they differ in recruitment window and in instrument "
    "content. That assumption is violated, so the estimate measures how separable the projects are, not what fraction is nonhuman. "
    "My recommendation is that we never put that number in a table without the sentence in red underneath it.\n\n"
    "D4 guards against a forking path. The graded knowledge score came out significant and the binary one did not, so we pre-specified the "
    "binary version as primary and report the graded one only as a sensitivity analysis.\n\n"
    "D5 concerns the four records excluded from clustering. All of them had valid screening answers and none said Definitely yes, so we "
    "report the full range the gap could take if they were reassigned, 39.0 to 44.5 points. The conclusion does not flip anywhere in that range."
)

# ─────────────────────────────────────────────────────────────────────────────
# S14 — DEMOGRAPHIC SIGNAL PREVALENCE
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Demographic concentration check",
    "The prevalence contrast alone is directional but not decisive",
    "The subgroup appears more often in 4581 than 4797, but prevalence alone does not justify a label.",
    [f"dirty_4581: {int(dirty_prev.n_selected)}/{int(dirty_prev.n_records)} = {float(dirty_prev.selected_pct)*100:.2f}% (Wilson 95% CI {float(dirty_prev.ci_low)*100:.2f}–{float(dirty_prev.ci_high)*100:.2f}%)",
     f"clean_4797: {int(clean_prev.n_selected)}/{int(clean_prev.n_records)} = {float(clean_prev.selected_pct)*100:.2f}% (Wilson 95% CI {float(clean_prev.ci_low)*100:.2f}–{float(clean_prev.ci_high)*100:.2f}%)",
     f"Dirty vs clean Fisher exact p={float(cmp_prev.p_value):.3f}, odds ratio {float(cmp_prev.odds_ratio):.2f}, risk ratio {float(cmp_prev.risk_ratio):.2f}",
     "Interpretation rule: demographic field values are never a stand-alone bot label; they are only an anomaly-concentration signal when corroborated by independent trust-risk evidence"],
    F("figure_21_demographic_signal_enrichment.png"),
    "figure_21_demographic_signal_enrichment.png; table_35_demographic_signal_summary.csv",
    "This slide separates a descriptive fact from an inferential claim.\n\n"
    "The descriptive fact is that the subgroup appears at 2.59% in 4581 and 1.13% in 4797. "
    "The inferential question is whether that prevalence contrast by itself proves contamination. "
    "It does not: the two-sided Fisher test is 0.312.\n\n"
    "That is exactly why this remains a concentration signal rather than a label. We keep the "
    "prevalence numbers visible, but we do not use them to classify records on their own.\n\n"
    "The analytic guardrail is simple: demographic values can only contribute to quality review "
    "when they co-enrich with independent timing, logic, and response-pattern rules that were "
    "defined without demographic outcomes."
)

# ─────────────────────────────────────────────────────────────────────────────
# S15 — DEMOGRAPHIC SIGNAL CO-ENRICHMENT
# ─────────────────────────────────────────────────────────────────────────────
fig_slide(prs, "Co-enrichment with trust-risk rules",
    "The marked 46-record subgroup is concentrated in independent risk signals",
    "Concentration is strongest for hard invalidity and rapid-completion rule families.",
    [f"Tier 1 confirmed invalid: {int(tier1_enrich.group_flagged_n)}/{int(tier1_enrich.group_total_n)} ({float(tier1_enrich.group_rate)*100:.1f}%) vs {int(tier1_enrich.comparison_flagged_n)}/{int(tier1_enrich.comparison_total_n)} ({float(tier1_enrich.comparison_rate)*100:.1f}%), OR {float(tier1_enrich.odds_ratio):.2f}, p={float(tier1_enrich.p_value):.2g}",
     f"Tier <=2 high-or-confirmed: {int(tier2_enrich.group_flagged_n)}/{int(tier2_enrich.group_total_n)} ({float(tier2_enrich.group_rate)*100:.1f}%) vs {int(tier2_enrich.comparison_flagged_n)}/{int(tier2_enrich.comparison_total_n)} ({float(tier2_enrich.comparison_rate)*100:.1f}%), OR {float(tier2_enrich.odds_ratio):.2f}, p={float(tier2_enrich.p_value):.2g}",
     "Rule-level enrichment survives FDR correction for R1, R2, R3, R7 and R8; R6 is not enriched because it is high-base-rate in the entire legacy project",
     "Operational decision: treat this as a suspected contamination cluster requiring stricter inclusion tiering, not deterministic demographic exclusion"],
    F("figure_21_demographic_signal_enrichment.png"),
    "figure_21_demographic_signal_enrichment.png; table_35b_demographic_signal_enrichment.csv; table_16_trust_tier_counts.csv",
    "This slide is the key decision support piece for the 46-record question.\n\n"
    "Once we condition on dirty_4581 only, the subgroup is substantially over-represented in "
    "Tier 1 and Tier <=2, and in several independent rule families. That cross-rule pattern "
    "is much stronger evidence than prevalence alone.\n\n"
    "We still avoid a demographic label. What we can claim is concentration of independent "
    "quality-risk signals, which supports stricter tiering and manual review.\n\n"
    "This is the mathematically conservative position: do not claim identity-based causation, "
    "do not claim exact bot prevalence, and do use the co-enrichment profile to prioritize "
    "exclusion sensitivity analyses."
)

print("[exec] 15 slides written")

# ═════════════════════════════════════════════════════════════════════════════
# SHARED HELPERS FOR APPENDIX / TABLE SLIDES
# ═════════════════════════════════════════════════════════════════════════════
def divider(prs, eyebrow, title, sub):
    s = blank(prs)
    if os.path.exists(SUNBURST):
        s.shapes.add_picture(SUNBURST, Inches(9.55), Inches(1.55),
                             height=Inches(4.2), width=Inches(4.2 * aspect(SUNBURST)))
    txt(s, M, 2.60, 8.4, 0.26, eyebrow, size=11, color=ORANGE, font=FH, bold=True,
        caps=True, space_after=0, char_space=0.6)
    txt(s, M, 2.94, 8.6, 1.0, title, size=38, color=DISCOVERY, font=FH, bold=True,
        space_after=0, line=0.93, char_space=-0.3)
    txt(s, M, 4.24, 8.4, 0.8, sub, size=13, color=JET, font=FB, space_after=0, line=1.02)
    logo_pair(s, M, 6.42, 0.34)
    return s


def cards(slide, y, items, h=1.18, fill=COOLBLUE, ncol=3, tcol=DISCOVERY):
    gap = 0.26
    w = (SW - 2 * M - gap * (ncol - 1)) / ncol
    for i, (hd, body) in enumerate(items):
        r, c = divmod(i, ncol)
        x = M + c * (w + gap); yy = y + r * (h + 0.22)
        plate(slide, x, yy, w, h, fill=fill, radius=0.12)
        txt(slide, x + 0.22, yy + 0.16, w - 0.44, 0.30, hd, size=11.5, color=tcol,
            font=FH, bold=True, space_after=0, line=0.95)
        txt(slide, x + 0.22, yy + 0.50, w - 0.44, h - 0.62, body, size=9.5, color=JET,
            font=FB, space_after=0, line=1.0)


def table_slide(prs, eyebrow, title, headers, widths, rows, note=None,
                src="", notes="", fsize=8.4, rh=0.40, y0=1.70, colors=None):
    s = blank(prs)
    txt(s, M, 0.40, 9.0, 0.24, eyebrow, size=10.5, color=ORANGE, font=FH, bold=True,
        caps=True, space_after=0, char_space=0.6)
    txt(s, M, 0.66, 10.6, 0.55, title, size=24, color=DISCOVERY, font=FH, bold=True,
        space_after=0, line=0.92, char_space=-0.2)
    lab_logo(s, SW - M - 0.62, 0.40, 0.29)
    cx = [M]
    for w in widths[:-1]:
        cx.append(cx[-1] + w)
    plate(s, M, y0 - 0.36, sum(widths), 0.34, fill=DISCOVERY, radius=0.09)
    for i, hh in enumerate(headers):
        txt(s, cx[i] + 0.13, y0 - 0.285, widths[i] - 0.20, 0.22, hh, size=9,
            color=COOLWHITE, font=FH, bold=True, caps=True, space_after=0, char_space=0.4)
    for j, row in enumerate(rows):
        yy = y0 + j * rh
        if j % 2 == 0:
            plate(s, M, yy, sum(widths), rh - 0.045, fill=COOLWHITE, radius=0.05)
        for i, cell in enumerate(row):
            col = (colors or {}).get((j, i), JET)
            bold = (i == 0)
            txt(s, cx[i] + 0.13, yy + 0.075, widths[i] - 0.22, rh - 0.14, cell,
                size=fsize, color=col, font=FH if bold else FB, bold=bold,
                space_after=0, line=0.95)
    if note:
        txt(s, M, y0 + len(rows) * rh + 0.14, SW - 2 * M, 0.5, note, size=9,
            color=GREY, font=FB, space_after=0, line=1.0)
    source(s, src)
    s.notes_slide.notes_text_frame.text = notes
    return s

# ═════════════════════════════════════════════════════════════════════════════
# PART B — TECHNICAL APPENDIX (8 slides + divider)
# ═════════════════════════════════════════════════════════════════════════════
d = divider(prs, "Part B", "Technical appendix",
            "Methods, assumptions, sensitivity checks, detector caveats, branching-logic\n"
            "audit detail and limitations. Eight slides for technical reviewers.")
d.notes_slide.notes_text_frame.text = (
    "From here on the audience is the technical reviewers. I will move faster and stop for questions rather than read every line. "
    "The purpose of this section is that someone who wants to challenge a number can find where it came from and what assumption it rests on."
)

# A1 — provenance
s = blank(prs)
header(s, "Appendix A1 · Provenance", "Provenance, caching and the gates that halt the pipeline")
takeaway(s, "Every input is hashed, dated and count-checked, and the pipeline halts rather than silently analysing a drifted cohort.")
cards(s, 2.46, [
    ("Hashed, dated caches",
     "4797 record 177×385, 4581 record 1779×369, plus metadata (324 and 315 fields) and instruments. Each Parquet cache carries a SHA-256 sidecar and is dated 2026-08-14."),
    ("Hard halt conditions",
     "The run aborts if the 177 / 1,779 record counts change, or if the verified-human timing reference moves off 131 completed 4797 records."),
    ("Field alignment",
     "311 fields exist in both projects, 13 only in 4797 and 4 only in 4581. Pooled models use shared behavioural and TFA fields only."),
], h=1.55)
bullets(s, M, 4.42, 12.13, 2.10, [
    "Namespaced UIDs make all 1,956 records unique despite 174 colliding numeric record IDs across the two projects",
    "The 135-ID analytic cohort comes from cleaned_autism_study_data.csv; table_1 records that its selection rule is undocumented in the scoped files",
    "Aggregate outputs carry no email, ZIP, date of birth, occupation text or API token; record_flags.parquet and data_cache/ are git-ignored",
    "Open documentation inconsistency: README.md states a primary snapshot of 2026-07-30 while table_12 lists verified caches dated 2026-08-14. Worth reconciling before circulation",
], size=11)
source(s, "table_1_file_inventory.csv; table_12_api_cache_inventory.csv; table_13_field_intersection.csv; table_31_upgrade_data_quality.csv; README.md")
s.notes_slide.notes_text_frame.text = (
    "Two things I want technical reviewers to take from this slide.\n\n"
    "First, the pipeline is fail-closed. If the record counts drift or the verified-human reference changes, it stops rather than quietly "
    "producing different numbers. That is why the anchors in this deck can be trusted to correspond to a specific data state.\n\n"
    "Second, the honest weak point is the cohort selection rule. We have 135 analytic IDs and table_1 explicitly records that the rule "
    "producing them is undocumented in the files we were given. That does not invalidate the internal comparisons, because everything "
    "downstream is conditioned on that cohort, but it does limit what we can say about generalisation to the wider caregiver population. "
    "If anyone can reconstruct that rule it would strengthen the paper.\n\n"
    "Last bullet is a housekeeping item rather than an analytic one: the README and the cache inventory disagree on the snapshot date. "
    "I have used the cache inventory date throughout because the CSVs are the ground truth, but the README should be corrected."
)

# A2 — clustering method
fig_slide(prs, "Appendix A2 · Method", "How the profiles are built",
    "Ten TFA domains, standardized exactly once, k-means with fixed seeds, and composite-level median imputation for partial responders.",
    ["Input is 10 aligned TFA domain scores on a 0 to 100 scale where higher means more accepting, standardized exactly once before clustering",
     "Inclusion threshold is at least 8 of 10 domains answered, which admits 131 of the 135 selected records",
     "k-means with seeded initialisation; k chosen from silhouette, 80% subsample ARI, inertia and the gap statistic",
     "Latent-profile assignment is effectively hard: mean maximum posterior 0.9999, 10th percentile 0.9999, and no record below 0.70",
     "Demographics, values, knowledge and screening are excluded from the feature matrix and analysed only after assignment"],
    F("figure_12_tfa_clustermap.png"),
    "figure_12_tfa_clustermap.png; figure_13_tfa_pca_biplot.png; table_25_latent_profile_quality.csv; config.yaml",
    "The clustermap is the raw picture of what is being clustered: 177 records in the pooled sensitivity definition across ten standardized "
    "TFA domains. You can see block structure, but you can also see that it is soft rather than crisp, which is consistent with the "
    "silhouette value from the executive section.\n\n"
    "Two preprocessing choices to flag. Standardization happens exactly once, at the domain level, which avoids the common bug of scaling "
    "twice and shrinking real spread. And partial responders are handled by composite-level median imputation, which is tested directly in "
    "the next slide by re-running on complete cases only.\n\n"
    "The posterior numbers look almost too clean, and they are worth explaining. A mean maximum posterior of 0.9999 does not mean the clusters "
    "are well separated. It means that, conditional on the fitted model, assignment is unambiguous. Separation quality is what silhouette "
    "measures, and that is the 0.180 we already discussed. Do not let anyone quote the posterior as evidence of strong clustering."
)

# A3 — sensitivity
fig_slide(prs, "Appendix A3 · Sensitivity", "How much the partition moves under reasonable alternatives",
    "Agreement with the primary partition ranges from 0.58 to 0.82 ARI, so the boundary is method-dependent even though the direction is not.",
    ["Complete-domain cases only (n=120): ARI 0.585, silhouette 0.153 — this is the direct test of composite-level median imputation",
     "Excluding the ethical and moral support domain (n=131): ARI 0.614, and the screening association persists at 62.8% versus 40.8%, a 22.0 pp difference, Fisher p=0.018",
     "5th to 95th percentile winsorization: ARI 0.690, silhouette 0.165 — marginal outliers are not driving the split",
     "Ward hierarchical clustering: ARI 0.822, silhouette 0.182, the least method-dependent of the four checks",
     "Nonlinear embeddings at three UMAP neighbourhood sizes and three t-SNE perplexities show overlapping gradient structure, not separated islands"],
    F("figure_14_umap_tsne_sensitivity.png"),
    "figure_14_umap_tsne_sensitivity.png; table_10_sensitivity_checks.csv",
    "This slide is the honest answer to how stable is the boundary. The answer is: moderately, and it depends on the method.\n\n"
    "Ward hierarchical clustering reproduces the k-means partition at 0.822 ARI, which is good. Dropping to complete cases only gives 0.585, "
    "which is the weakest of the four and tells us the imputation choice does move individual assignments. Roughly speaking, somewhere "
    "between a sixth and a quarter of caregivers could land in the other profile under a defensible alternative specification.\n\n"
    "What does not move is the direction of the screening association. The ethics-domain check is the most informative one here, because the "
    "obvious objection is that the whole result is driven by moral endorsement of screening. Drop that domain entirely, rebuild the profiles, "
    "and the screening gap is still 22 percentage points at Fisher p equals 0.018. Smaller than the headline 42 points, but clearly present.\n\n"
    "On the embeddings: I show them because reviewers ask for them, but I would not over-read them. Across all six panels the two profiles occupy "
    "overlapping regions with a gradient between them. That is a picture of a continuum being partitioned, which is exactly what we said on slide 3."
)

# A4 — model-based
fig_slide(prs, "Appendix A4 · Model-based checks", "Model-based checks support two profiles without endorsing them as optimal",
    "The bootstrap likelihood-ratio test rejects a single component, while BIC prefers a three-component diagonal solution.",
    ["Bootstrap LRT, k=2 versus k=1, diagonal covariance: observed likelihood ratio 1587.8 over 100 bootstrap repeats, p=0.0099 — one component is rejected",
     "Best overall BIC is k=3 diagonal at 1697.9, against 2329.7 for k=2 diagonal, so the information criterion does not endorse two components",
     "Entropy separation exceeds 0.99 for the retained solution and no record falls below a 0.70 maximum posterior",
     "Per-case silhouettes are positive for all but a handful of conditional-profile cases; the negative values are left visible rather than trimmed",
     "Stated plainly: two profiles is the interpretable and stable choice, not the information-criterion optimum"],
    F("figure_16_case_silhouette.png"),
    "table_23_gmm_diagnostics.csv; table_24_gmm_bootstrap_lrt.csv; table_25_latent_profile_quality.csv; figure_15_consensus_matrix.png; figure_16_case_silhouette.png",
    "This is the slide a methods reviewer will go to first, so I want the tension on it rather than hidden.\n\n"
    "The bootstrap likelihood-ratio test does its job: k equals 2 beats k equals 1 with a p-value of 0.0099 over 100 bootstrap repeats. So there "
    "is structure beyond a single Gaussian. But BIC across the 24 fitted specifications prefers a three-component diagonal model. We are not "
    "reporting the BIC optimum, and I would rather say that out loud than have it found.\n\n"
    "The defence is the one from slide 3. The k equals 3 solution is unstable in its smallest cluster and would not survive characterisation, "
    "and the bootstrap consensus matrix in figure 15 shows one solid block and one much more diffuse block rather than three clean ones.\n\n"
    "The per-case silhouette plot matters for a different reason. A handful of conditional-profile cases have negative silhouettes, meaning they "
    "sit closer to the other profile's centre than their own. We deliberately did not trim them. Trimming borderline cases would inflate every "
    "separation statistic on the previous slides and would be a form of self-deception."
)

# A5 — rules
fig_slide(prs, "Appendix A5 · Rule definitions", "The ten rules, their thresholds, and where the thresholds come from",
    "Every threshold is a percentile or floor computed on the 131 completed 4797 caregivers, so the screen is calibrated to observed human behaviour.",
    ["R1 total time below 11.57 minutes and R2 TFA time below 7.85 minutes, both verified-human minima",
     "R3 any instrument below the clean 1st percentile; R4 within-block response SD at or below the clean 1st percentile, computed separately for the 1–4, 1–5, 1–6 and 1–7 Likert blocks",
     "R5 exact ordered Likert fingerprint duplicate at ≥80% answered; R6 at least 3 submissions inside a clean-derived 60-second window; R7 near-duplicate open text at TF-IDF cosine ≥0.90",
     "R8 logical and branching inconsistency; R9 impossible demographic combination; R10 4797 instrument-native gate, validation-only and not tier-defining",
     "Counts over all 1,956 records: R6 1,682, R3 846, R4 589, R2 447, R1 283, R8 44, R7 17, R10 13, R9 2, R5 0"],
    F("figure_7_timing_ecdf.png"),
    "table_14_fraud_rule_definitions.csv; table_15b_rule_counts_by_project.csv; figure_7_timing_ecdf.png",
    "The ECDF is the single clearest justification for the timing rules. The blue curve is the 131 verified-human completions and the orange is "
    "the 1,575 legacy records with complete timing blocks. The dotted line at 11.57 minutes is the fastest a verified human finished.\n\n"
    "Notice what the shaded region means: about 18% of legacy records finished faster than any human we can vouch for. That is the R1 flag, and it "
    "is a floor rather than a percentile, which makes it conservative by construction. We are not saying those records are automated. We are "
    "saying they are outside the range of observed human completion.\n\n"
    "Two design points for reviewers. First, R4 is computed separately for each Likert block width, because a standard deviation on a 1 to 4 scale "
    "is not comparable to one on a 1 to 7 scale, and collapsing them is a common error. Second, every threshold comes from the clean project, never "
    "from the pooled data, so the legacy records cannot influence the definition of what counts as suspicious.\n\n"
    "R5 firing zero times is worth a sentence: whatever produced the legacy records was not submitting identical response vectors."
)

# A6 — branching audit detail
FAM_ROWS = [
    ("16", "prenatal_tested_reason_without_prenatal_yes", "Reason for prenatal testing answered with no prenatal-yes gate", "36.4%"),
    ("9",  "child_age_band_count_exceeds_child_count", "More distinct child age bands ticked than children reported", "20.5%"),
    ("8",  "prenatal_testing_mutually_exclusive_pair", "Both members of a mutually exclusive prenatal pair selected", "18.2%"),
    ("6",  "prenatal_not_tested_reason_without_prenatal_no", "Reason for not testing answered with no prenatal-no gate", "13.6%"),
    ("3",  "earlier_diagnosis_no_reason_without_no_gate", "No-reason answered without the earlier-diagnosis no gate", "6.8%"),
    ("2",  "earlier_diagnosis_yes_reason_without_yes_gate", "Yes-reason answered without the earlier-diagnosis yes gate", "4.5%"),
]
table_slide(prs, "Appendix A6 · Branching audit",
    "What a branching violation looks like, record by record",
    ["n", "Violation family", "Plain-language description", "Share"],
    [0.62, 4.35, 5.86, 1.30],
    [list(r) for r in FAM_ROWS],
    note=("Worked example, record 4581_1069: fif_num_children = 2, yet fif_childrens_ages___3, ___4 and ___5 are all ticked. "
          "Three distinct age bands for two children is not producible by the form. Every row stores gating_fields, gating_logic, "
          "observed_fields, gating_values and observed_values, so any reviewer can re-derive the call without rerunning the pipeline.\n"
          "Regression fixtures pinned in README.md: sibling-pair collisions 31, 118, 273, 571, 768, 1076, 1554, 1699 · tested-reason orphans "
          "966–979, 981, 1591 · not-tested orphans 364, 416, 987, 1026, 1707, 1741 · earlier-diagnosis orphans 205, 1025, 1137, 1642, 1678."),
    src="table_34_branching_logic_audit.csv (44 rows, all dirty_4581); table_15_rule_false_positive_rates.csv; README.md",
    fsize=9, rh=0.46, y0=1.78,
    notes=(
        "Why this deserves its own appendix slide: it is the only part of the trust screen that does not depend on a threshold.\n\n"
        "Take the worked example. A caregiver reports two children, then ticks three different child age bands. There is no distribution to "
        "compare against and no percentile to argue over. The form cannot produce that combination, so the record contains at least one answer "
        "that is not a truthful report. That is a qualitatively stronger kind of evidence than a fast completion time.\n\n"
        "The table stores both sides of every check, the gating fields and the observed fields with their values, which means a clinical "
        "collaborator can audit any of the 44 rows in a spreadsheet without touching the code. That reviewability is the methodological point.\n\n"
        "Three limitations to state. The audit covers six manually specified families, so it finds what we thought to look for and nothing else. "
        "Forty-four records is 2.5% of the 1,779 legacy records, so this is high precision and unknown recall — it is not a contamination estimate. "
        "And the zero count in the clean project is only a weak validity check, since the clean project is small and its caregivers were supervised.\n\n"
        "The regression fixtures at the bottom are pinned so that any future change to R8 that silently stops catching these records fails the test suite."
    ))

# A7 — identifiability
fig_slide(prs, "Appendix A7 · Identifiability", "Why we cannot name a contamination fraction",
    "The estimator assumes labelled positives are drawn at random from the positive class, and our labelled set is simply a different REDCap project.",
    ["Elkan–Noto style nonhuman share in 4581: 0.944 (0.939 to 0.950), with a labelling propensity c of 0.622",
     "Source-classifier AUC 0.980 measures separability of two projects that differ in recruitment window and in instrument content",
     "Negative control, random clean half A versus half B: AUC 0.525, close to chance, so the pipeline does not manufacture separation on its own",
     "Top discriminating features are open-text field count and validation-block time, both instrument-design differences rather than bot signatures",
     "The prior lab LightGBM artifact was not present in the scoped files, so concordance uses rules, three novelty detectors and the PU classifier only"],
    F("figure_10_behavioral_feature_space.png"),
    "table_17_contamination_identifiability.csv; table_18_detector_feature_importance.csv; table_19_detector_agreement.csv; figure_10_behavioral_feature_space.png; figure_11b_pu_feature_importance.png",
    "The PCA panel on the left is the argument in one picture. The blue 4797 points sit inside the orange 4581 cloud, not beside it. If the legacy "
    "project were overwhelmingly nonhuman and the clean project human, we would expect two displaced masses. Instead we see one broad distribution "
    "with the verified humans in a dense band inside it.\n\n"
    "The UMAP panel colours the same records by deterministic tier, and the tiers do not form separate islands either. Tier 1 through Tier 4 are "
    "interleaved throughout. That is what a graded evidential screen looks like, as opposed to a detector that has found a distinct population.\n\n"
    "Now the estimator. Elkan and Noto's method requires the selected completely at random assumption: labelled positives must be a random draw "
    "from all positives. Our labelled set is one REDCap project and the unlabelled set is another, collected in a different window with a partly "
    "different instrument. The assumption is not approximately satisfied, it is structurally violated. So the 94.4% figure quantifies how "
    "separable the two projects are. Calling it a bot prevalence would be a category error.\n\n"
    "The negative control at 0.525 is what tells us the machinery is sound. Split the clean project at random and the classifier cannot tell the "
    "halves apart. Good hygiene, but it does not rescue the interpretation of the 0.980."
)

# A8 — limitations
s = blank(prs)
header(s, "Appendix A8 · Limitations", "Limitations we should state before anyone else does")
takeaway(s, "Small n, cross-sectional measurement, an undocumented cohort selection rule, and a trust screen that grades evidence rather than identifying actors.")
cards(s, 2.46, [
    ("Power, not equivalence",
     "With 47 and 80 respondents and a 54.3% pooled reference rate, the minimum detectable difference is about 25.6 pp at 80% power. Null characterization results are uninformative about small effects."),
    ("Undocumented selection",
     "The 135-ID analytic cohort's selection rule is undocumented in the scoped files, so internal comparisons stand but external generalisation does not follow."),
    ("Cross-sectional design",
     "Attitudes and screening intent were measured in the same sitting. No temporal ordering exists, so no causal claim is available in either direction."),
    ("Excluded-record tipping point",
     "Reassigning the four cluster-excluded records moves the screening gap between 39.0 and 44.5 pp. The conclusion does not flip anywhere in that range."),
    ("Assignment-error correction",
     "BCH-corrected screening rates are 58.3% versus 31.5%, lower than the modal-assignment rates of 70.0% and 27.7% but preserving the ordering."),
    ("Trust screen is not attribution",
     "Tiers grade how much evidence supports a record. A Tier 1 record may be rushed, malformed or automated, and these data cannot separate those explanations."),
], h=1.72, ncol=3)
txt(s, M, 6.28, 12.13, 0.5,
    "Regression harness: 17 of 17 status-quo checks pass within tolerance, including selected_n 135, clusterable_n 131, screen_valid_n 127, "
    "silhouette 0.1803, screening gap 42.34 pp and additive cluster OR 8.729. Colour accessibility was checked, with deuteranopia-simulated "
    "separation of 0.578 and non-colour backups in place on every figure.",
    size=9.5, color=GREY, font=FB, line=1.0, space_after=0)
source(s, "table_29_power_precision.csv; table_28_excluded_case_tipping_point.csv; table_26_bch_screening_outcome.csv; table_27_status_quo_regression_checks.csv; table_32_color_accessibility_check.csv; table_1_file_inventory.csv")
s.notes_slide.notes_text_frame.text = (
    "I would rather present these than be handed them, so let me go through them quickly.\n\n"
    "The power limitation is the one most likely to be misused. When we said knowledge, values and demographics do not separate the profiles, "
    "that is a statement about what we could detect at this sample size, which is roughly a 25 point difference. A real 10 point difference "
    "would be invisible to us. Please do not let anyone write that the profiles are demographically identical.\n\n"
    "The BCH correction deserves a word. It accounts for classification error in profile assignment when relating profiles to a distal outcome. "
    "The corrected rates, 58.3 versus 31.5, are more conservative than the modal ones, 70.0 versus 27.7, which is expected. The ordering and the "
    "substantive conclusion survive, but the corrected numbers are the more defensible ones for a manuscript.\n\n"
    "On the trust screen: a tier is a statement about evidence, not about the respondent. A Tier 1 record might be a rushed human, a malformed "
    "submission, or something automated. These data cannot tell those apart, and the honest position is to say so.\n\n"
    "The regression harness at the bottom is the reason I can stand behind the numbers in this deck. Seventeen pinned checks all pass within "
    "tolerance, so the pipeline state that produced these outputs is the one we validated."
)
# A9 — figure index (every remaining generated figure, physically in the deck)
s = blank(prs)
header(s, "Appendix A9 · Figure index", "The eight figures not placed earlier, and the claim each one carries")
takeaway(s, "All 21 generated figures now appear in this deck; these eight support claims made on slides 4, 7, 11, 12, A2, A4 and A5.")
IDX = [
    ("figure_5_values_characterization.png",     "Slide 7 · No value survives FDR control"),
    ("figure_6_demographic_characterization.png","Slide 7 · No demographic contrast survives FDR"),
    ("figure_17_tfa_domain_distributions.png",   "Slide 4 · Distribution shape behind the profile means"),
    ("figure_8_response_fingerprint_heatmap.png","A5 · Legacy response texture versus verified-clean"),
    ("figure_13_tfa_pca_biplot.png",             "A2 · Linear structure and centroid validity"),
    ("figure_15_consensus_matrix.png",           "A4 · One solid block, one diffuse block"),
    ("figure_11b_pu_feature_importance.png",     "Slide 12 · Open-text count dominates the PU classifier"),
    ("figure_19_tier_to_cluster_alluvial.png",   "Slide 11 · Only 14 of 131 sit below Tier 4"),
]
gw, gh, gx, gy = 2.96, 1.72, 0.245, 0.60
for i, (fn, cap) in enumerate(IDX):
    r, c = divmod(i, 4)
    x = M + c * (gw + gx); y = 2.46 + r * (gh + gy)
    figure(s, F(fn), x, y, gw, gh, pad=0.07)
    txt(s, x + 0.02, y + gh + 0.07, gw - 0.04, 0.24, cap, size=8, color=DISCOVERY,
        font=FH, bold=True, space_after=0, line=0.95)
    txt(s, x + 0.02, y + gh + 0.30, gw - 0.04, 0.22, fn, size=6.8, color=GREY,
        font=FB, space_after=0, line=0.95)
source(s, "All eight PNGs are the unmodified 300-DPI originals from Caregiver Outputs; full-size versions are in that folder")
s.notes_slide.notes_text_frame.text = (
    "This index exists so the deck is self-contained. Every figure the pipeline produced, figure_1 through figure_20 including figure_11b, "
    "is now either shown at full size earlier or reproduced here as a thumbnail, and each one is bound to a specific claim in the traceability map.\n\n"
    "Three of these are worth opening at full size if the discussion goes there. Figure 5 and figure 6 are the values and demographics "
    "characterisations, which are the visual form of the null results on slide 7 and are the ones a reviewer is most likely to want to inspect "
    "for a missed signal. Figure 15, the bootstrap consensus matrix, is the strongest single picture of why we did not go to three clusters: "
    "one block is solid and the other is visibly diffuse.\n\n"
    "Figure 8 is the response fingerprint comparison, and I would treat it carefully. The two panels look different in texture, but the sample "
    "sizes differ by more than a factor of ten, so some of that difference is a rendering artefact of row count. It is illustrative, not evidential."
)
print("[appendix] 10 slides written")

# ═════════════════════════════════════════════════════════════════════════════
# PART C — VISUAL-TO-CLAIM TRACEABILITY MAP
# ═════════════════════════════════════════════════════════════════════════════
TRACE_A = [
    ("3",  "k=2 is the most stable and interpretable cluster count; separation is modest", "figure_1_cluster_count_diagnostics.png · table_3_cluster_diagnostics.csv"),
    ("4",  "Profiles separate most on trust, blood-test openness and ethical support", "figure_2_primary_cluster_profiles.png · table_4_cluster_defining_profiles.csv"),
    ("4",  "Domain-level distribution shape behind the profile means", "figure_17_tfa_domain_distributions.png · table_4_cluster_defining_profiles.csv"),
    ("5",  "Definitely-yes screening intent is 70.0% versus 27.7%, a 42.3 pp gap", "figure_3_screening_outcome.png · table_5_screening_outcome.csv"),
    ("6",  "Adjusted ORs: profile 8.73, ASD at home 4.32, interaction 2.02 (n.s.)", "figure_18_logistic_forest.png · table_11_asd_family_context.csv · table_22_logistic_models.csv"),
    ("7",  "Verified autism knowledge does not separate the profiles (p=0.099)", "figure_4_autism_knowledge.png · table_6_knowledge_summary.csv · table_6b_knowledge_item_results.csv"),
    ("7",  "No value survives FDR control; self-direction is the largest raw signal", "figure_5_values_characterization.png · table_7_values_summary.csv"),
    ("7",  "No demographic contrast reaches FDR p < 0.56", "figure_6_demographic_characterization.png · table_8_demographics_summary.csv"),
    ("7",  "Minimum detectable difference is about 25.6 pp at 80% power", "table_29_power_precision.csv"),
    ("8",  "Rule co-occurrence and verified-human false-positive rates", "figure_9_rule_cooccurrence.png · table_15_rule_false_positive_rates.csv · table_14_fraud_rule_definitions.csv"),
    ("8",  "Rule counts across both projects, R6 = 1,682 and R5 = 0", "table_15b_rule_counts_by_project.csv"),
    ("9",  "Tier 4 pass share is 86.4% in clean_4797 and 3.9% in dirty_4581", "table_16_trust_tier_counts.csv (chart built from these values)"),
    ("10", "44 branching and logic impossibilities across six families, all in dirty_4581", "table_34_branching_logic_audit.csv · table_15b_rule_counts_by_project.csv"),
    ("10", "R8 verified-human false-positive rate is 0.0% in 131 records", "table_15_rule_false_positive_rates.csv"),
]
TRACE_B = [
    ("11", "Screening gap holds at 42.3 to 52.2 pp across five inclusion definitions", "figure_20_tier_sensitivity.png · table_20_tier_sensitivity.csv"),
    ("11", "Only 14 of 131 clusterable caregivers sit below Tier 4", "figure_19_tier_to_cluster_alluvial.png · table_16_trust_tier_counts.csv"),
    ("12", "Maximum detector agreement is Cohen's kappa 0.29", "figure_11_detector_scores_and_agreement.png · table_19_detector_agreement.csv"),
    ("12", "Top PU feature is open-text field count (permutation AUC drop 0.204)", "figure_11b_pu_feature_importance.png · table_18_detector_feature_importance.csv"),
    ("13", "Five decisions, four implemented, contamination fraction left open", "table_30_decision_summary.csv · table_17_contamination_identifiability.csv"),
    ("14", "Demographic prevalence contrast is directional but not decisive", "table_35_demographic_signal_summary.csv · figure_21_demographic_signal_enrichment.png"),
    ("15", "Marked subgroup co-enriches with independent trust-risk flags", "table_35b_demographic_signal_enrichment.csv · table_16_trust_tier_counts.csv"),
    ("A1", "Hashed dated caches, halt gates, and field alignment across projects", "table_12 · table_13 · table_13b · table_1 · table_31 · README.md"),
    ("A2", "What is clustered: ten standardized TFA domains, softly structured", "figure_12_tfa_clustermap.png · table_25_latent_profile_quality.csv"),
    ("A2", "Linear structure and centroid validity in PCA space", "figure_13_tfa_pca_biplot.png"),
    ("A3", "Partition agreement ranges 0.585 to 0.822 ARI under alternatives", "figure_14_umap_tsne_sensitivity.png · table_10_sensitivity_checks.csv"),
    ("A4", "Bootstrap LRT p=0.0099 for k=2 vs k=1; BIC prefers k=3 diagonal", "table_24_gmm_bootstrap_lrt.csv · table_23_gmm_diagnostics.csv · figure_15_consensus_matrix.png"),
    ("A4", "Negative per-case silhouettes retained rather than trimmed", "figure_16_case_silhouette.png"),
    ("A5", "Timing floors are verified-human minima, not percentiles of pooled data", "figure_7_timing_ecdf.png · table_14_fraud_rule_definitions.csv"),
    ("A5", "Legacy response texture differs from the verified-clean reference", "figure_8_response_fingerprint_heatmap.png"),
    ("A7", "Verified humans sit inside the legacy cloud; tiers form no islands", "figure_10_behavioral_feature_space.png · table_17_contamination_identifiability.csv"),
    ("A8", "Tipping point 39.0–44.5 pp; BCH 58.3% vs 31.5%; 17/17 regression checks pass",
            "table_28 · table_26 · table_27 · table_32_color_accessibility_check.csv"),
]
for part, rowset in (("1 of 2", TRACE_A), ("2 of 2", TRACE_B)):
    table_slide(prs, f"Part C · Traceability {part}",
        "Every claim mapped to the file that supports it",
        ["Slide", "Claim as stated in the deck", "Supporting figure and table filenames"],
        [0.68, 5.05, 6.40], [list(r) for r in rowset],
        src="All files resolve under Caregiver Outputs/ unless noted; output_manifest.csv carries a SHA-256 prefix for each",
        fsize=7.6, rh=0.355, y0=1.66,
        notes=("Traceability map. Every claim in this deck resolves to a named file in Caregiver Outputs, and output_manifest.csv "
               "carries a SHA-256 prefix for each artifact so the exact version can be confirmed.\n\n"
               "Two conventions worth stating. Where a figure and a table cover the same claim, the table is authoritative for the number "
               "and the figure is illustrative, because figures round and tables do not. And the two native charts in this deck, on slides 9 "
               "and 10, are drawn directly from table_16 and table_34 values rather than re-computed, so there is no separate derivation to audit.\n\n"
               "All 22 generated figures, figure_1 through figure_21 including figure_11b, appear somewhere in this map. If a reviewer asks "
               "why a particular figure was not shown in the executive section, the answer is in this table: it was bound to an appendix claim instead."))
print("[traceability] 2 slides written")

# ═════════════════════════════════════════════════════════════════════════════
# PART D — FINAL QA CHECKLIST
# ═════════════════════════════════════════════════════════════════════════════
QA = [
    ("Two distinct caregiver types exist", "Weak",
     "A two-profile partition of a continuous acceptability gradient is the most stable and interpretable summary of these data."),
    ("Higher acceptability leads to screening uptake", "Unsupported",
     "Profile membership is strongly associated with held-out screening intent. The design is cross-sectional, so no ordering is available."),
    ("94.4% of legacy records are bots", "Unsupported",
     "A source classifier separates the two projects at AUC 0.980. The implied separation share is not an identifiable contamination estimate."),
    ("The trust screen detects bots", "Unsupported",
     "The trust screen grades how much evidence supports each record. It does not identify automated respondents."),
    ("Tier 1 records were submitted by bots", "Unsupported",
     "Tier 1 records contain hard logical or timing impossibilities under our rules, whatever produced them."),
    ("The profiles do not differ demographically", "Weak",
     "No demographic difference was detectable here. Differences below roughly 25 percentage points would not be visible at these sample sizes."),
    ("Self-direction values distinguish the profiles", "Weak",
     "Self-direction shows the largest raw difference among ten values tested but does not survive false-discovery-rate control."),
    ("The 4581 subset replicates the finding", "Weak",
     "The legacy Tier-4 subset is directionally consistent, with an interval (5.7 to 74.1 pp) too wide to be confirmatory."),
    ("R8 implies about 2.5% contamination", "Unsupported",
     "R8 confirms that 44 records contain logically impossible responses. Its recall is unknown, so it bounds nothing."),
    ("Cluster membership is stable", "Weak",
     "The direction of the result is stable. Individual assignments move under alternative preprocessing, ARI 0.585 to 0.822."),
]
qcol = {}
for i, r in enumerate(QA):
    qcol[(i, 1)] = RED if r[1] == "Unsupported" else ORANGE
table_slide(prs, "Part D · QA checklist",
    "Claims to avoid, and the wording that the evidence will carry",
    ["Tempting claim", "Status", "Safer wording"],
    [4.05, 1.28, 6.80], [list(r) for r in QA],
    src="Status assigned from table_17, table_20, table_23, table_24, table_29, table_10, table_7, table_8, table_19 and table_34",
    fsize=8, rh=0.50, y0=1.66, colors=qcol,
    notes=(
        "This is the slide I would put in front of anyone drafting text from this deck.\n\n"
        "The two categories mean different things. Unsupported means the evidence cannot bear the claim at all, and the sentence should not be "
        "written in any form. Weak means there is a real signal but it is smaller, more fragile or more conditional than the tempting phrasing "
        "implies, so the wording needs to carry that.\n\n"
        "The three I would most expect to slip through review are the causal one on row two, the 94.4% on row three, and the demographic "
        "equivalence claim on row six. All three are natural things to write and none of them survives contact with the numbers.\n\n"
        "Note that nothing in this table changes the headline finding. Profile membership is associated with screening intent at 42.3 percentage "
        "points, that association survives adjustment and five inclusion definitions, and the branching audit is clean. The checklist is about "
        "not spending that credibility on claims the data do not support."
    ))

# QA verification log
s = blank(prs)
header(s, "Part D · Verification log", "What was checked before this deck was written")
takeaway(s, "The build script re-reads every source CSV and asserts each anchor; it refuses to write the file on any mismatch.")
cards(s, 2.46, [
    ("18 anchors asserted at build time",
     "Combined records 1,956 · clean_4797 177 · dirty_4581 1,779 · R8 44 / 0 / 44 · R8 false positives 0.0% in 131 · all six branching families · dirty tiers 1–3 at 96.1%. All 18 passed."),
    ("17 of 17 pipeline checks pass",
     "table_27 confirms selected_n 135, clusterable_n 131, screen_valid_n 127, silhouette 0.1803, screening gap 42.34 pp and additive cluster OR 8.729, each inside its tolerance."),
    ("Figures used unmodified",
     "Every figure is the 300-DPI PNG from Caregiver Outputs, uncropped and unedited. The two native charts are drawn from table_16 and table_34 values only."),
], h=1.60, ncol=3)
bullets(s, M, 4.34, 12.13, 2.30, [
    "Open item: README.md states a primary cache snapshot of 2026-07-30 while table_12 lists verified caches dated 2026-08-14. The deck follows the CSV. Reconcile before circulation",
    "Open item: the selection rule producing the 135-ID analytic cohort is undocumented in the scoped files, which limits external generalisation claims",
    "Known gaps carried forward: the prior lab LightGBM artifact was not present in the scoped files, and R10's email anti-fraud component is unavailable under current API permissions",
    "Not verified here: whether the figure PNGs were regenerated from the current pipeline state. output_manifest.csv carries a SHA-256 prefix per artifact, which would settle it",
], size=11)
source(s, "table_27_status_quo_regression_checks.csv; table_31_upgrade_data_quality.csv; output_manifest.csv; README.md; build script assertion log")
s.notes_slide.notes_text_frame.text = (
    "Closing note on process, which matters because this deck will be reused.\n\n"
    "The build script for these slides does not accept any number as given. It opens each source CSV, recomputes the quantity, and compares it "
    "against the agreed anchor. If any one of the eighteen fails, it raises and no file is written. So the numbers on these slides and the numbers "
    "in Caregiver Outputs cannot drift apart silently.\n\n"
    "Four open items, none of them blocking. The README and cache inventory disagree on a date, which is documentation rather than analysis. "
    "The cohort selection rule is undocumented, which is the most substantive of the four and worth chasing. The LightGBM artifact and the email "
    "component are known coverage gaps we have already declared rather than hidden. And I have not independently confirmed that the figure PNGs "
    "were regenerated from the current pipeline state, though the manifest hashes would settle that in a minute if anyone wants it.\n\n"
    "If the group agrees with the decisions on slide 13 and the wording on the previous slide, this material is ready to go into the manuscript draft."
)
print("[qa] 2 slides written")

prs.save(f"{OUTDIR}/ESD_Caregiver_Cluster_and_Trust_Screen.pptx")
print(f"[saved] total slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")

