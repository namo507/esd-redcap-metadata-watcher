"""ESD Lab deck construction helpers (brand-locked)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

SK = "/sessions/trusting-cool-heisenberg/mnt/.claude/skills/esd-lab"
OUT = "/sessions/trusting-cool-heisenberg/mnt/caregiver-cluster-analysis/Caregiver Outputs"

# --- canon brand tokens (non-negotiable) ---
DISCOVERY = RGBColor(0x33, 0x66, 0xFF)
SCIENCE   = RGBColor(0x91, 0xBA, 0xF4)
COOLBLUE  = RGBColor(0xE6, 0xEE, 0xFC)
COOLWHITE = RGBColor(0xF4, 0xF4, 0xF6)
JET       = RGBColor(0x00, 0x00, 0x00)
ORANGE    = RGBColor(0xF5, 0x7F, 0x00)
RED       = RGBColor(0xD7, 0x4E, 0x2D)
YELLOW    = RGBColor(0xF4, 0xDA, 0x26)
PINK      = RGBColor(0xF8, 0xB2, 0xB1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x5A, 0x5A, 0x64)

FH = "Libre Franklin"          # headings, bold=True
FB = "Libre Franklin Medium"   # body

SW, SH = 13.333, 7.5
M = 0.6

LOGO_LAB_BLUE  = f"{SK}/assets/logos/logo-horizontal-discovery-blue.png"
LOGO_LAB_WHITE = f"{SK}/assets/logos/logo-horizontal-cool-white.png"
LOGO_UOFSC     = f"{SK}/assets/logos/uofsc-horizontal-garnet.png"
PATTERN_BAND   = f"{SK}/assets/patterns/pattern-icon-band-white.png"
SUNBURST       = f"{SK}/assets/icons/sunburst-discovery-blue.png"
STAR_ORANGE    = f"{SK}/assets/icons/star-confident-orange.png"


def aspect(path):
    with Image.open(path) as im:
        return im.size[0] / im.size[1]


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def blank(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def txt(slide, x, y, w, h, runs, size=12, color=JET, font=FB, bold=False,
        align=PP_ALIGN.LEFT, space_after=6, line=None, anchor=MSO_ANCHOR.TOP,
        caps=False, char_space=None):
    """runs: str or list of str (one paragraph each)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tb.text_frame.vertical_anchor = anchor
    items = [runs] if isinstance(runs, str) else runs
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line:
            p.line_spacing = line
        r = p.add_run()
        r.text = t.upper() if caps else t
        f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
        if char_space is not None:
            from pptx.oxml.ns import qn
            f._rPr.set('spc', str(int(char_space * 100)))
    return tb


def bullets(slide, x, y, w, h, items, size=12, color=JET, gap=9, line=0.98):
    """Bulleted list with a true hanging indent so wrapped lines align under the text."""
    hang = Inches(0.055 * size / 3.0)          # scales with type size
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = line
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(hang)))
        pPr.set("indent", str(-int(hang)))
        rm = p.add_run(); rm.text = "–\t"
        rm.font.name = FB; rm.font.size = Pt(size); rm.font.color.rgb = DISCOVERY; rm.font.bold = True
        r = p.add_run(); r.text = t
        r.font.name = FB; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def plate(slide, x, y, w, h, fill=COOLWHITE, radius=0.035, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.adjustments[0] = radius
    return sh


def lab_logo(slide, x, y, h, white=False):
    p = LOGO_LAB_WHITE if white else LOGO_LAB_BLUE
    a = aspect(p)
    slide.shapes.add_picture(p, Inches(x), Inches(y), height=Inches(h), width=Inches(h * a))
    return h * a


def logo_pair(slide, x, y, h, white=False):
    """Lab logo LEFT of UofSC logo, both same height (brand rule)."""
    w1 = lab_logo(slide, x, y, h, white=white)
    gap = 0.22
    a2 = aspect(LOGO_UOFSC)
    slide.shapes.add_picture(LOGO_UOFSC, Inches(x + w1 + gap), Inches(y),
                             height=Inches(h), width=Inches(h * a2))
    return w1 + gap + h * a2


def header(slide, eyebrow, title, title_color=DISCOVERY, title_size=26):
    txt(slide, M, 0.40, 9.0, 0.24, eyebrow, size=10.5, color=ORANGE, font=FH,
        bold=True, caps=True, space_after=0, char_space=0.6)
    txt(slide, M, 0.66, 9.45, 0.85, title, size=title_size, color=title_color,
        font=FH, bold=True, space_after=0, line=0.92, char_space=-0.2)
    lab_logo(slide, SW - M - 0.62, 0.40, 0.29)


def takeaway(slide, text, y=1.62, h=0.60, w=SW - 2 * M):
    plate(slide, M, y, w, h, fill=COOLBLUE, radius=0.30)
    txt(slide, M + 0.24, y + 0.055, w - 0.48, h - 0.11, text, size=12.5,
        color=DISCOVERY, font=FH, bold=True, space_after=0, line=0.95,
        anchor=MSO_ANCHOR.MIDDLE)


def source(slide, text):
    txt(slide, M, SH - 0.44, SW - 2 * M, 0.24, "Source: " + text, size=8,
        color=GREY, font=FB, space_after=0)


def figure(slide, path, x, y, w, h, pad=0.10):
    """Fit image inside box (x,y,w,h) on a cool-white rounded card."""
    a = aspect(path)
    iw, ih = w - 2 * pad, (w - 2 * pad) / a
    if ih > h - 2 * pad:
        ih = h - 2 * pad; iw = ih * a
    cx = x + (w - iw) / 2 - pad
    cy = y + (h - ih) / 2 - pad
    plate(slide, cx, cy, iw + 2 * pad, ih + 2 * pad, fill=COOLWHITE, radius=0.05)
    slide.shapes.add_picture(path, Inches(cx + pad), Inches(cy + pad),
                             width=Inches(iw), height=Inches(ih))


def fig_slide(prs, eyebrow, title, take, bl, fig, src, notes, force=None,
              bullet_size=12):
    """Auto layout: 'split' (bullets left / figure right) or 'wide'."""
    s = blank(prs)
    header(s, eyebrow, title)
    takeaway(s, take)
    a = aspect(fig)
    mode = force or ("wide" if a >= 2.05 else "split")
    if mode == "split":
        bullets(s, M, 2.48, 4.55, 4.30, bl, size=bullet_size)
        figure(s, fig, 5.42, 2.42, 7.30, 4.45)
    else:
        half = (SW - 2 * M - 0.40) / 2
        n = len(bl); k = (n + 1) // 2
        bullets(s, M, 2.44, half, 1.35, bl[:k], size=10.5, gap=6)
        bullets(s, M + half + 0.40, 2.44, half, 1.35, bl[k:], size=10.5, gap=6)
        figure(s, fig, M, 3.92, SW - 2 * M, 2.98)
    source(s, src)
    s.notes_slide.notes_text_frame.text = notes
    return s
