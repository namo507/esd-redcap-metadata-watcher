#!/usr/bin/env python3
"""Brand QA for the built deck.

Checks the things the esd-lab skill calls non-negotiable, on the saved file
rather than on the build script's intentions:

  1. every colour on a slide is a canon hex, and neither drifted blue appears
  2. the theme palette and typeface are branded, so a later edit in PowerPoint
     cannot pick Office blue out of the colour picker
  3. every text run is Libre Franklin, and every character it contains is one
     Libre Franklin can actually draw (no silent font fallback)
  4. no placeholder or template residue survived
  5. the lab logo sits left of the UofSC logo at the same height, everywhere

    python3 deck/qa.py [file.pptx]
"""

from __future__ import annotations

import os
import re
import sys
import zipfile

from fontTools.ttLib import TTFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = os.path.dirname(os.path.abspath(__file__))
DEFAULT_ASSETS = (
    "/Users/namomac/Library/Application Support/Claude/local-agent-mode-sessions/"
    "skills-plugin/c898a6b1-6fa4-467a-8bf4-94014b33c02c/"
    "07ab60b6-de84-467a-ae49-82bf671a695e/skills/esd-lab/assets"
)
FONT_DIR = os.path.join(os.environ.get("ESD_ASSETS", DEFAULT_ASSETS), "fonts")

CANON = {
    "3366FF", "91BAF4", "E6EEFC", "F4F4F6", "000000", "F57F00", "D74E2D",
    "F4DA26", "F8B2B1", "FFFFFF",
    "1A1A1F", "4A4A55",  # ink and muted, both derived neutrals
}
DRIFT = {"005CBE", "2A61E6"}
BRAND_FONTS = {"Libre Franklin", "Libre Franklin Medium"}
PLACEHOLDERS = ("lorem", "ipsum", "todo", "[insert", "click to edit")


def main() -> int:
    path = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        os.path.dirname(ROOT), "ESD-Visit-Scheduling-v3.pptx"
    )
    fails: list = []
    notes: list = []

    zf = zipfile.ZipFile(path)

    # 1. slide colours only (the theme is checked separately)
    slide_hexes = set()
    for name in zf.namelist():
        if not re.match(r"ppt/slides/slide\d+\.xml$", name):
            continue
        xml = zf.read(name).decode("utf-8", "ignore")
        slide_hexes |= {h.upper() for h in re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', xml)}
    if DRIFT & slide_hexes:
        fails.append(f"drifted blue on a slide: {sorted(DRIFT & slide_hexes)}")
    non_canon = slide_hexes - CANON
    if non_canon:
        fails.append(f"non-canon hexes on slides: {sorted(non_canon)}")
    notes.append(f"slide colours: {len(slide_hexes)} distinct, all canon"
                 if not non_canon else "")

    # 2. theme
    theme_name = next(n for n in zf.namelist() if n.startswith("ppt/theme/theme"))
    theme = zf.read(theme_name).decode("utf-8", "ignore")
    accent1 = re.search(r"<a:accent1><a:srgbClr val=\"([0-9A-Fa-f]{6})\"", theme)
    if not accent1 or accent1.group(1).upper() != "3366FF":
        fails.append("theme accent1 is not discovery blue")
    if "Libre Franklin" not in theme:
        fails.append("theme typeface is not Libre Franklin")

    # 3. fonts and glyph coverage
    coverage = set()
    for filename in ("LibreFranklin-Regular.ttf", "LibreFranklin-Medium.ttf",
                     "LibreFranklin-Bold.ttf"):
        for table in TTFont(os.path.join(FONT_DIR, filename))["cmap"].tables:
            coverage |= set(table.cmap.keys())

    prs = Presentation(path)
    fonts_used, missing = set(), {}
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts_used.add(run.font.name)
                    for ch in run.text:
                        if ch.strip() and ord(ch) not in coverage:
                            missing.setdefault(ch, set()).add(i)
    if not fonts_used <= BRAND_FONTS:
        fails.append(f"non-brand fonts: {sorted(fonts_used - BRAND_FONTS)}")
    if missing:
        fails.append(
            "characters Libre Franklin cannot draw (PowerPoint would silently "
            "fall back): "
            + ", ".join(f"{c!r} on slides {sorted(v)}" for c, v in missing.items())
        )

    # 4. placeholders
    all_text = " ".join(
        sh.text_frame.text for sl in prs.slides for sh in sl.shapes
        if sh.has_text_frame
    ).lower()
    for word in PLACEHOLDERS:
        if word in all_text:
            fails.append(f"placeholder text present: {word!r}")

    # 5. logo pairing
    for i, slide in enumerate(prs.slides, start=1):
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        lab = [p for p in pics if "logo-horizontal" in (p.image.filename or "")]
        usc = [p for p in pics if "uofsc" in (p.image.filename or "")]
        if lab and usc:
            if lab[0].left >= usc[0].left:
                fails.append(f"slide {i}: lab logo is not left of the UofSC logo")
            if abs(lab[0].height - usc[0].height) > 1000:
                fails.append(f"slide {i}: logo heights differ")
        elif usc and not lab:
            fails.append(f"slide {i}: UofSC logo without the lab logo")
        elif not pics:
            fails.append(f"slide {i}: no logo at all")

    print(f"file      {os.path.basename(path)}")
    print(f"slides    {len(prs.slides._sldIdLst)}")
    print(f"colours   {' '.join(sorted(slide_hexes))}")
    print(f"fonts     {', '.join(sorted(fonts_used))}")
    print(f"theme     accent1 {accent1.group(1) if accent1 else '?'}, "
          f"typeface {'Libre Franklin' if 'Libre Franklin' in theme else '?'}")
    print()
    if fails:
        print("QA FAILURES")
        for f in fails:
            print(f"  - {f}")
        return 1
    print("BRAND QA PASSED")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
