#!/usr/bin/env python3
"""Build the ESD Lab visit-scheduling v3 deck.

    python3 deck/render_math_v3.py      # writes deck/math/*.png + math_dims.json
    python3 deck/build_deck_v3.py       # writes ESD-Visit-Scheduling-v3.pptx

Brand tokens, type ramp and slide recipes follow the esd-lab skill (Path B,
built fresh rather than from the template, because none of the template's 32
layouts carries a formula band). pptxgenjs is not installed on this machine, so
this is the python-pptx equivalent of the skill's Path B recipes.

Environment:
    ESD_ASSETS  the esd-lab skill's assets/ directory
"""

from __future__ import annotations

import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.abspath(__file__))
MATH = os.path.join(ROOT, "math")
DIMS_PATH = os.path.join(ROOT, "math_dims.json")
OUT = os.path.join(
    os.path.dirname(ROOT), "ESD-Visit-Scheduling-v3.pptx"
)

DEFAULT_ASSETS = (
    "/Users/namomac/Library/Application Support/Claude/local-agent-mode-sessions/"
    "skills-plugin/c898a6b1-6fa4-467a-8bf4-94014b33c02c/"
    "07ab60b6-de84-467a-ae49-82bf671a695e/skills/esd-lab/assets"
)
ASSETS = os.environ.get("ESD_ASSETS", DEFAULT_ASSETS)

# ---------------------------------------------------------------------------
# Brand tokens (canon hexes, no drift)
# ---------------------------------------------------------------------------

C = {
    "discovery": RGBColor(0x33, 0x66, 0xFF),
    "science": RGBColor(0x91, 0xBA, 0xF4),
    "coolBlue": RGBColor(0xE6, 0xEE, 0xFC),
    "coolWhite": RGBColor(0xF4, 0xF4, 0xF6),
    "jet": RGBColor(0x00, 0x00, 0x00),
    "orange": RGBColor(0xF5, 0x7F, 0x00),
    "red": RGBColor(0xD7, 0x4E, 0x2D),
    "yellow": RGBColor(0xF4, 0xDA, 0x26),
    "pink": RGBColor(0xF8, 0xB2, 0xB1),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "ink": RGBColor(0x1A, 0x1A, 0x1F),
    "muted": RGBColor(0x4A, 0x4A, 0x55),
}
FH = "Libre Franklin"          # headings, bold
FB = "Libre Franklin Medium"   # body

W, H = 13.333, 7.5
M = 0.62  # page margin

DIMS = json.load(open(DIMS_PATH, encoding="utf-8"))


def asset(*parts: str) -> str:
    return os.path.join(ASSETS, *parts)


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
prs.core_properties.title = "Visit Scheduling Scoring System v3"
prs.core_properties.author = "The Early Social Development Lab at UofSC"

BLANK = prs.slide_layouts[6]


def slide(bg=None):
    s = prs.slides.add_slide(BLANK)
    if bg is not None:
        rect(s, 0, 0, W, H, bg)
    return s


def rect(s, x, y, w, h, fill, radius=None, line=None, line_w=1.0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        # adjustment is a fraction of the shorter side
        sh.adjustments[0] = min(0.5, radius / min(w, h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def text(
    s, string, x, y, w, h,
    size=14, color=C["ink"], bold=False, font=None, align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP, spacing=1.15, caps=False, char_spacing=None,
):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = string.split("\n") if isinstance(string, str) else list(string)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line.upper() if caps else line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font or (FH if bold else FB)
        r.font.color.rgb = color
        if char_spacing is not None:
            r.font._rPr.set("spc", str(int(char_spacing * 100)))
    return box


def eyebrow(s, string, x, y, color=None):
    return text(
        s, string, x, y, 8, 0.28, size=11, bold=True,
        color=color or C["discovery"], caps=True, char_spacing=1.2,
    )


def math_img(s, key, y, x=None, cx=None, height=None, max_w=None):
    """Place a rendered formula at its natural aspect ratio."""
    pw, ph = DIMS[key]
    h = height if height is not None else ph / 400.0 * 0.92
    w = h * pw / ph
    if max_w and w > max_w:
        w = max_w
        h = w * ph / pw
    if cx is not None:
        x = cx - w / 2
    s.shapes.add_picture(os.path.join(MATH, f"{key}.png"), Inches(x), Inches(y),
                         Inches(w), Inches(h))
    return w, h


def logo_pair(s, variant="discovery-blue", y=0.32, height=0.40, x=None):
    """Lab logo LEFT of the UofSC logo, same height. Never reordered."""
    lab = asset("logos", f"logo-horizontal-{variant}.png")
    from PIL import Image

    lw, lh = Image.open(lab).size
    lab_w = height * lw / lh
    uofsc = asset("logos", "uofsc-horizontal-garnet.png")
    uw, uh = Image.open(uofsc).size
    u_w = height * uw / uh
    x = M if x is None else x
    s.shapes.add_picture(lab, Inches(x), Inches(y), Inches(lab_w), Inches(height))
    s.shapes.add_picture(
        uofsc, Inches(x + lab_w + 0.30), Inches(y), Inches(u_w), Inches(height)
    )


def header(s, kicker, title, title_color=None, y=1.30):
    # Logo pair occupies 0.32 to 0.72 in; the eyebrow has to clear it.
    eyebrow(s, kicker, M, y - 0.38)
    text(s, title, M, y, W - 2 * M, 0.8, size=32, bold=True,
         color=title_color or C["discovery"], char_spacing=-0.3)
    return y + 0.74


LOCAL_ICONS = os.path.join(ROOT, "icons")


def icon(s, name, x, y, size=0.72, folder="suite"):
    # Recoloured variants live beside the deck; the science-blue originals are
    # invisible on a science-blue tile, and recolour.py is the sanctioned way
    # to restate single-hue art in another canon colour.
    local = os.path.join(LOCAL_ICONS, name)
    path = local if os.path.exists(local) else (
        asset("icons", folder, name) if folder else asset("icons", name))
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(size), Inches(size))


def bullets(s, items, x, y, w, size=14, gap=0.44, color=None, bullet_color=None):
    """Small square bullets rather than glyph bullets, so spacing is controllable."""
    cy = y
    for item in items:
        rect(s, x, cy + 0.085, 0.11, 0.11, bullet_color or C["discovery"], radius=0.05)
        lines = 1 + len(item) // 66
        text(s, item, x + 0.30, cy, w - 0.30, 0.34 * lines, size=size,
             color=color or C["muted"], spacing=1.25)
        cy += gap + 0.20 * (lines - 1)
    return cy


def table(s, headers, rows, x, y, w, col_w=None, row_h=0.40, size=12,
          header_fill=None, zebra=True, align=None):
    n = len(headers)
    col_w = col_w or [w / n] * n
    align = align or [PP_ALIGN.LEFT] * n
    rect(s, x, y, w, row_h, header_fill or C["discovery"])
    cx = x
    for i, head in enumerate(headers):
        text(s, head, cx + 0.14, y + 0.085, col_w[i] - 0.2, row_h,
             size=size - 2.5, bold=True, color=C["white"], caps=True,
             char_spacing=0.8, align=align[i])
        cx += col_w[i]
    cy = y + row_h
    for j, row in enumerate(rows):
        if zebra and j % 2 == 0:
            rect(s, x, cy, w, row_h, C["coolWhite"])
        cx = x
        for i, cell in enumerate(row):
            bold = isinstance(cell, tuple)
            value, colour = (cell if bold else (cell, None))
            text(s, str(value), cx + 0.14, cy + 0.095, col_w[i] - 0.2, row_h,
                 size=size, color=colour or C["ink"], bold=bold,
                 align=align[i], font=FH if bold else FB)
            cx += col_w[i]
        cy += row_h
    return cy


def stat_tile(s, x, y, w, h, value, label, fill=None, value_color=None):
    rect(s, x, y, w, h, fill or C["discovery"], radius=0.18)
    text(s, value, x + 0.28, y + 0.20, w - 0.56, 0.75, size=40, bold=True,
         color=value_color or C["white"], char_spacing=-0.5)
    text(s, label, x + 0.28, y + 0.98, w - 0.56, h - 1.05, size=12,
         color=value_color or C["white"], spacing=1.2)


def footer(s, n, dark=False):
    text(s, f"{n:02d}", W - M - 0.6, H - 0.55, 0.6, 0.3, size=11, bold=True,
         color=C["science"] if dark else C["muted"], align=PP_ALIGN.RIGHT)


# ===========================================================================
# Slides
# ===========================================================================

n = 0


def new(kicker=None, title=None, bg=None, title_color=None, dark=False):
    global n
    n += 1
    s = slide(bg)
    if kicker or title:
        logo_pair(s, "cool-white" if dark else "discovery-blue")
        y = header(s, kicker or "", title or "", title_color)
    else:
        y = 1.0
    footer(s, n, dark)
    return s, y


# --- 1. title --------------------------------------------------------------
n += 1
s = slide(C["discovery"])
s.shapes.add_picture(asset("patterns", "pattern-icon-band-white.png"),
                     Inches(0), Inches(0.0), Inches(W), Inches(0.92))
s.shapes.add_picture(asset("patterns", "pattern-icon-band-white.png"),
                     Inches(0), Inches(H - 0.92), Inches(W), Inches(0.92))
logo_pair(s, "cool-white", y=1.55, height=0.5)
text(s, "Visit Scheduling\nScoring System", M, 2.45, 9.2, 2.0, size=54, bold=True,
     color=C["coolWhite"], char_spacing=-0.6, spacing=1.02)
rect(s, M, 4.72, 1.5, 0.055, C["yellow"])
text(s, "Version 3  ·  fair and efficient coordinator assignment for home visits",
     M, 5.0, 9.5, 0.4, size=17, color=C["coolWhite"])
text(s, "Pilot build  ·  17 August 2026  ·  engine v3.0.0  ·  weights provisional",
     M, 5.55, 9.5, 0.35, size=13, color=C["coolWhite"])
s.shapes.add_picture(asset("icons", "sunburst-cool-white.png"),
                     Inches(W - M - 1.9), Inches(2.6), Inches(1.9), Inches(1.9))

# --- 2. what the notes asked ----------------------------------------------
s, y = new("From the meeting notes", "Four next steps, four implementations")
rows = [
    ("Calendar data integration via Graph / API",
     "Graph getSchedule + Google freeBusy, staleness policy, write-time recheck"),
    ("Travel converted to a workload metric",
     "One prospective burden term, with the exchange rate made explicit"),
    ("Family history + recency merged for continuity",
     "One continuity index, saturating in visits and decaying in days"),
    ("Debrief unexpected scoring outcomes",
     "Named detector rules, weekly report generated from the audit log"),
]
table(s, ["Meeting note", "Where it landed in v3"], rows, M, y + 0.05,
      W - 2 * M, col_w=[4.6, W - 2 * M - 4.6], row_h=0.66, size=13)
text(s, "Everything below follows from taking those four seriously, plus three "
        "findings that came out of checking the v2 arithmetic.",
     M, y + 3.25, W - 2 * M, 0.5, size=14, color=C["muted"])

# --- 3. the three findings -------------------------------------------------
s, y = new("Why v3 and not v2.1", "Three findings from checking the arithmetic")
cw = (W - 2 * M - 0.6) / 3
for i, (num, head, body, colour) in enumerate([
    ("1", "Travel was priced at 13x",
     "Under min-max normalisation the 0.20 / 0.15 split implied one hour of "
     "driving was worth about thirteen hours of clinic work. Nobody would "
     "defend that out loud.", C["red"]),
    ("2", "Pool-relative scoring caused cold start",
     "Normalising to the busiest teammate makes scores non-stationary and hands "
     "a brand new coordinator a 1.0. The missing history was never the bug.",
     C["orange"]),
    ("3", "0.65 of the weight sat on one construct",
     "Family history, recency and the continuity bonus are three slots measuring "
     "relationship state once. Splitting a construct inflates its total.",
     C["discovery"]),
]):
    x = M + i * (cw + 0.3)
    rect(s, x, y, cw, 3.35, C["coolBlue"], radius=0.2)
    rect(s, x, y, 0.09, 3.35, colour)
    text(s, num, x + 0.35, y + 0.28, 1.0, 0.6, size=34, bold=True, color=colour)
    text(s, head, x + 0.35, y + 1.0, cw - 0.7, 0.9, size=17, bold=True,
         color=C["discovery"], spacing=1.1)
    text(s, body, x + 0.35, y + 1.95, cw - 0.7, 1.3, size=12.5, color=C["muted"])
math_img(s, "gamma_implied", y + 3.62, cx=W / 2, height=0.62)

# --- 4. architecture -------------------------------------------------------
s, y = new("End to end", "Four layers, and nothing skips a layer")
lw = (W - 2 * M - 0.75) / 4
layers = [
    ("Layer 0", "Freshness gate", "How old is the calendar? fresh / stale / expired",
     "icon-checklist-cool-white.png"),
    ("Layer 1", "Hard eligibility", "Seven predicates, ANDed. No score can rescue a failure",
     "icon-hands-pair-discovery-blue.png"),
    ("Layer 2", "Weighted score", "Four non-redundant criteria, each in [0, 1]",
     "icon-growth-chart-cool-white.png"),
    ("Layer 3", "Rank + review", "Top 3, calibrated band, stability, logged ties",
     "icon-eye-discovery-blue.png"),
]
for i, (tag, head, body, ic) in enumerate(layers):
    x = M + i * (lw + 0.25)
    rect(s, x, y, lw, 2.55, C["discovery"] if i % 2 == 0 else C["science"], radius=0.2)
    icon(s, ic, x + 0.3, y + 0.28, 0.62)
    text(s, tag, x + 0.3, y + 1.02, lw - 0.6, 0.28, size=11, bold=True,
         color=C["coolWhite"] if i % 2 == 0 else C["ink"], caps=True, char_spacing=1.0)
    text(s, head, x + 0.3, y + 1.32, lw - 0.6, 0.4, size=17, bold=True,
         color=C["white"] if i % 2 == 0 else C["ink"])
    text(s, body, x + 0.3, y + 1.78, lw - 0.6, 0.7, size=11.5,
         color=C["coolWhite"] if i % 2 == 0 else C["muted"], spacing=1.2)
    if i < 3:
        text(s, "›", x + lw + 0.02, y + 1.05, 0.25, 0.4, size=24, bold=True,
             color=C["discovery"], align=PP_ALIGN.CENTER)
math_img(s, "composite", y + 2.95, cx=W / 2, height=0.5)
text(s, "Feasible pool empty? The visit is returned for manual scheduling with the "
        "reason each coordinator failed. It is never rescued by a high score.",
     M, y + 3.72, W - 2 * M, 0.5, size=13, color=C["muted"], align=PP_ALIGN.CENTER)

# --- 5. Layer 0 ------------------------------------------------------------
s, y = new("Layer 0", "Calendar staleness is a status, never a score penalty")
math_img(s, "layer1", y - 0.06, x=M, height=0.42)
rows = [
    ("Visit within 72 hours", "15 min", "60 min"),
    ("3 to 14 days out", "4 hours", "24 hours"),
    ("More than 14 days out", "24 hours", "72 hours"),
]
table(s, ["Visit horizon", "Hard: block", "Soft: flag"], rows, M, y + 0.72, 6.1,
      col_w=[3.0, 1.55, 1.55], row_h=0.44,
      align=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])
bullets(s, [
    "fresh: score and auto-commit",
    "stale: score normally, mark the assignment PROVISIONAL, hold the family "
    "notification until a human confirms",
    "expired or sync failed: Layer 1 failure",
    "Write-time recheck of the winning pair before commit. One API call, and it "
    "closes the fifteen minute race that double-books someone",
    "Circuit breaker: halt the run if more than 20% of the team is unverifiable",
], 6.9, y + 0.72, W - M - 6.9, size=13, gap=0.40)
rect(s, M, y + 2.55, 6.1, 1.15, C["coolBlue"], radius=0.15)
text(s, "A penalty would make a stale but ideal candidate lose for a reason that "
        "has nothing to do with fit, and quietly corrupt the metric every other "
        "diagnostic is built on.",
     M + 0.3, y + 2.75, 5.5, 0.8, size=12.5, color=C["muted"], spacing=1.25)

# --- 6. Layer 1 ------------------------------------------------------------
s, y = new("Layer 1", "Seven predicates. A failure is never outscored")
rows = [
    ("W", "Date window match", "Visit window overlaps declared working hours"),
    ("A", "Open slot", "Free block of duration PLUS round-trip travel"),
    ("X", "No calendar clash", "Nothing hard-booked over it: busy, out of office"),
    ("E", "No family conflict", "Hard exclusion list. Checked first, most sensitive"),
    ("K", "Credential match", "Req(protocol) is a subset of Cred(coordinator)"),
    ("Cal", "Calendar fresh", "Layer 0 gate"),
    ("Ramp", "Onboarding cap", "New hires capped at q visits per week"),
]
table(s, ["", "Predicate", "What it checks"], rows, M, y, W - 2 * M,
      col_w=[0.8, 2.9, W - 2 * M - 3.7], row_h=0.42, size=13)
text(s, "Slots are travel-inflated. Without that the engine happily schedules two "
        "visits forty minutes apart across town.",
     M, y + 3.34, W - 2 * M, 0.4, size=13, color=C["muted"])
rect(s, M, y + 3.85, W - 2 * M, 0.85, C["coolBlue"], radius=0.15)
text(s, "Ramp is a constraint, not a score adjustment. “Do not overload a new "
        "hire in week one” is a policy, and policies belong where they cannot be "
        "traded away against a good continuity score.",
     M + 0.3, y + 4.02, W - 2 * M - 0.6, 0.6, size=13, color=C["discovery"], bold=True)

# --- 7. Layer 2 overview ---------------------------------------------------
s, y = new("Layer 2", "Five terms became four criteria")
rows = [
    ("v2: family history 0.30", "-->", ("Phi  continuity  0.45", C["discovery"])),
    ("v2: recency 0.25", "-->", ""),
    ("v2: continuity bonus 0.10", "-->", ("P  protocol continuity  0.10", C["discovery"])),
    ("v2: workload 0.20", "-->", ("Psi  burden relief  0.30", C["discovery"])),
    ("v2: travel 0.15", "-->", ""),
    ("(not in v2 as a term)", "-->", ("Omega  family preference  0.15", C["discovery"])),
]
table(s, ["v2 term", "", "v3 criterion"], rows, M, y, 7.3,
      col_w=[3.4, 0.6, 3.3], row_h=0.40, size=12.5,
      align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT])
bullets(s, [
    "Family history keeps the highest weight, per the notes",
    "Family preference becomes its own criterion, per the notes",
    "Four criteria need six pairwise comparisons, not ten. That is a twenty "
    "minute elicitation meeting instead of an afternoon",
    "These weights are analyst-assigned and must not survive the pilot "
    "unvalidated",
], 8.0, y + 0.12, W - M - 8.0, size=13, gap=0.56)
math_img(s, "composite", y + 3.20, cx=W / 2, height=0.55)

# --- 8. Phi ----------------------------------------------------------------
s, y = new("Criterion Phi", "Continuity: saturating in visits, decaying in days")
math_img(s, "phi_raw", y, x=M, height=0.62)
math_img(s, "phi_flip", y + 0.95, x=M, height=0.85)
rows = [
    ("0", "0.00", "never met"),
    ("1", "0.39", "one visit"),
    ("3", "0.78", "familiar"),
    ("5", "0.92", "saturated"),
]
table(s, ["k prior visits", "familiarity", ""], rows, 7.6, y, W - M - 7.6,
      col_w=[1.9, 1.5, W - M - 7.6 - 3.4], row_h=0.40, size=12.5,
      align=[PP_ALIGN.CENTER, PP_ALIGN.RIGHT, PP_ALIGN.LEFT])
text(s, "kappa = 2, tau = 75 days.   sigma = +1 the family wants a familiar face, "
        "−1 a fresh one.",
     7.6, y + 2.05, W - M - 7.6, 0.6, size=12.5, color=C["muted"])
rect(s, M, y + 2.35, W - 2 * M, 1.5, C["coolBlue"], radius=0.18)
math_img(s, "phi_zero", y + 2.58, x=M + 0.4, height=0.42)
text(s, "Cold start dies by construction, not by imputation. Days-since-contact is "
        "UNDEFINED when they have never met, and it is never evaluated, because the "
        "product has already zeroed. The v2 bug was treating undefined as extreme.",
     M + 0.4, y + 3.12, W - 2 * M - 0.8, 0.7, size=13, color=C["muted"], spacing=1.25)

# --- 9. Psi + the repricing -----------------------------------------------
s, y = new("Criterion Psi", "Burden relief, and the price of merging travel in")
math_img(s, "psi_burden", y, x=M, height=0.68)
math_img(s, "psi_relief", y + 1.0, x=M, height=0.72)
bullets(s, [
    "Prospective: this visit's own duration and trip are in the number",
    "Capacity-referenced, so part-time coordinators work natively",
    "gamma is elicited in one question: how many extra minutes of clinic time "
    "would you accept to avoid ten minutes of driving?",
], M, y + 2.0, 6.2, size=13, gap=0.42)
rows = [
    ("1  pure time", "0.019", "÷ 7.8"),
    ("2  default", "0.037", "÷ 4.1"),
    ("3", "0.053", "÷ 2.8"),
    ("12.9", ("0.150", C["red"]), "÷ 1.0"),
]
table(s, ["gamma", "effective travel weight", "vs v2's 0.15"], rows, 7.0, y + 0.15,
      W - M - 7.0, col_w=[1.6, 2.6, W - M - 7.0 - 4.2], row_h=0.42, size=12.5,
      align=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])
text(s, "The merge is a repricing, and the lab has to see the price. If the team "
        "still wants travel to matter more, that is about EQUITY of travel, not "
        "burden, and equity belongs in the constraints.",
     7.0, y + 2.35, W - M - 7.0, 1.1, size=12.5, color=C["muted"], spacing=1.25)

# --- 10. cold start --------------------------------------------------------
s, y = new("Cold start", "The team median is wrong for two data types in three")
rows = [
    ("Observed and true", "scheduled hours = 0",
     "No shrinkage. An empty calendar really is empty"),
    ("Undefined by construction", "days since contact, k = 0",
     "Define, do not impute. R = 0 via the product"),
    ("Estimated with noise", "visit-duration multiplier",
     "Shrinkage belongs here, and only here"),
]
table(s, ["Data type", "Example", "Correct handling"], rows, M, y, W - 2 * M,
      col_w=[3.0, 3.4, W - 2 * M - 6.4], row_h=0.52, size=13)
math_img(s, "ramp", y + 2.1, x=M, height=0.72)
math_img(s, "shrinkage", y + 3.0, x=M, height=0.62)
bullets(s, [
    "The median cliff at N_min re-ranks discontinuously on the N_min-th visit",
    "It is also gameable: log a trivial visit to cross the threshold",
    "Method of moments and the precision criterion both give N_min ≈ 4m",
    "With 4 to 8 coordinators, do NOT fit m from pilot data: it is badly "
    "estimated and can go negative. Use m = 5, revisit at ~200 visits",
], 7.4, y + 2.1, W - M - 7.4, size=12.5, gap=0.40)

# --- 11. reference case ----------------------------------------------------
s, y = new("Correctness anchor", "Same three coordinators. v2 said C, A, B")
rows = [
    ("A", "3", "5 d", "18 h", "30 min", "yes", ("0.502", C["discovery"]), "1st"),
    ("B", "0", "—", "9 h", "75 min", "no", "0.173", "3rd"),
    ("C", "1", "40 d", "12 h", "45 min", "no", "0.246", "2nd"),
]
table(s, ["", "k", "since", "booked", "travel", "prev cp", "v3 score", "v3 rank"],
      rows, M, y, W - 2 * M,
      col_w=[0.7, 0.8, 1.1, 1.2, 1.3, 1.3, 2.0, W - 2 * M - 8.4], row_h=0.46, size=13,
      align=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 5 + [PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])
stat_tile(s, M, y + 2.35, 3.4, 1.55, "A, C, B", "v3 ranking, asserted to three "
          "decimals in the test suite")
stat_tile(s, M + 3.7, y + 2.35, 3.4, 1.55, "C, A, B",
          "v2 ranking on identical inputs", fill=C["coolBlue"],
          value_color=C["discovery"])
rect(s, M, y + 4.22, W - 2 * M, 0.82, C["coolBlue"], radius=0.15)
text(s, "The v2 worked example and the v3 test suite use the same table on "
        "purpose, so the two systems can be diffed rather than argued about.",
     M + 0.3, y + 4.43, W - 2 * M - 0.6, 0.5, size=13, bold=True,
     color=C["discovery"])
text(s, "Why it flipped: v2's pool-relative workload handed A a zero and C a free "
        "half point, and v2's recency term rewarded B simply for having been idle. "
        "A is genuinely at capacity in v3 too, but continuity plus same-rater "
        "protocol continuity carry the decision.",
     M + 7.4, y + 2.35, W - M - (M + 7.4), 1.6, size=13, color=C["muted"], spacing=1.3)

# --- 12. Layer 3 -----------------------------------------------------------
s, y = new("Layer 3", "The review band should be measured, not guessed")
math_img(s, "epsilon", y + 0.14, x=M, height=0.44, max_w=W - 2 * M)
stat_tile(s, M, y + 1.02, 2.9, 1.62, "0.020", "calibrated band, pilot data")
stat_tile(s, M + 3.15, y + 1.02, 2.9, 1.62, "0.050", "the v2 guess",
          fill=C["coolBlue"], value_color=C["discovery"])
text(s, "The guess was twice as conservative as the weights warrant, so it was "
        "sending twice as many decisions to a human as it needed to. A real pilot "
        "could move it either way. That is the point of measuring it.",
     M + 6.3, y + 1.02, W - M - (M + 6.3), 1.62, size=13, color=C["muted"], spacing=1.3)
bullets(s, [
    "Selection stability: P(top-1) over the weight simplex. Winning 55% of the "
    "mass is not the same decision as winning 99% at the same gap",
    "Stability under 0.60 routes to a human automatically",
    "Tie-break order: protocol continuity, then family-history direction, then "
    "uniform random WITH THE SEED LOGGED",
    "A randomised true tie is defensible on equipoise grounds, and turns every "
    "tie into a free experiment",
], M, y + 2.72, W - 2 * M, size=13, gap=0.44)

# --- 13. optimisation ------------------------------------------------------
s, y = new("Batch optimisation", "Two nested problems, and a trigger for each")
text(s, "A.  One coordinator, one day", M, y, 5.8, 0.32, size=15, bold=True,
     color=C["discovery"])
math_img(s, "dp", y + 0.38, x=M, height=0.42, max_w=5.9)
text(s, "Weighted interval scheduling. Sort by end time, exact, O(n log n), about "
        "fifteen lines. Intervals are travel-inflated.",
     M, y + 0.95, 5.9, 0.6, size=12.5, color=C["muted"])
math_img(s, "dp_trigger", y + 1.62, x=M, height=0.52)
text(s, "Conflicts appear at the fourth candidate per coordinator per day, or about "
        "ten open visits per coordinator per week. ESD is already there.",
     M, y + 2.3, 5.9, 0.6, size=12.5, color=C["muted"])

text(s, "B.  Whole team, whole week", 7.0, y, 5.8, 0.32, size=15, bold=True,
     color=C["discovery"])
math_img(s, "assignment", y + 0.38, x=7.0, height=0.5, max_w=5.9)
math_img(s, "assignment_st", y + 1.0, x=7.0, height=0.4, max_w=5.9)
rows = [
    ("Several visits per person", "no", "YES"),
    ("Part-time capacity", "no", "YES"),
    ("Leave a visit unfilled", "no", "YES"),
    ("Intra-day overlap", "no", "no"),
]
table(s, ["", "Hungarian", "Min-cost flow"], rows, 7.0, y + 1.55, 5.72,
      col_w=[2.9, 1.4, 1.42], row_h=0.36, size=11.5,
      align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
text(s, "Hungarian is the wrong shape: it is the square, one-visit-per-person case. "
        "Flow cannot express intra-day overlap, which IS problem A, so the two are "
        "solved together in a Benders-style repair loop.",
     7.0, y + 3.58, 5.9, 0.9, size=12.5, color=C["muted"], spacing=1.25)

# --- 14. escalation --------------------------------------------------------
s, y = new("Escalation", "The optimiser runs in shadow from day one")
math_img(s, "regret", y, x=M, height=0.85)
stat_tile(s, 5.6, y - 0.05, 2.35, 1.7, "8.8%", "regret on the pilot week",
          fill=C["orange"])
stat_tile(s, 8.2, y - 0.05, 2.35, 1.7, "3%", "escalation threshold",
          fill=C["coolBlue"], value_color=C["discovery"])
stat_tile(s, 10.8, y - 0.05, 2.0, 1.7, "0", "unfilled gap",
          fill=C["coolBlue"], value_color=C["discovery"])
bullets(s, [
    "Shadow mode computes, logs to optimizer_shadow, and changes nothing",
    "Escalate when regret exceeds 3% or a visit is lost, two weeks running",
    "Free intermediate fix, already in the production greedy path: process "
    "visits MOST-CONSTRAINED-FIRST. One line. Greedy's worst failure is burning "
    "the only ADOS-certified coordinator on a visit any generalist could take",
], M, y + 2.05, W - 2 * M, size=13, gap=0.44)
rect(s, M, y + 3.55, W - 2 * M, 0.75, C["coolBlue"], radius=0.15)
text(s, "The synthetic pilot week already clears the bar. On real data, two "
        "consecutive weeks over 3% is the signal to promote the optimiser.",
     M + 0.3, y + 3.72, W - 2 * M - 0.6, 0.5, size=13, bold=True, color=C["discovery"])

# --- 15. fairness ----------------------------------------------------------
s, y = new("Fairness", "Equity is a constraint, not a criterion")
text(s, "Encoding “nobody should always draw the long drives” as a weighted "
        "criterion makes it tradeable against everything else, which is exactly what "
        "fairness is meant to prevent.",
     M, y, W - 2 * M, 0.7, size=15, color=C["ink"], spacing=1.3)
cw = (W - 2 * M - 0.6) / 3
for i, (head, body, ic) in enumerate([
    ("Vetoes, in the engine",
     "Reject if the assignment pushes utilisation over 1.0, or if the rolling "
     "four-week travel share exceeds 1.4x the capacity share.",
     "icon-hands-pair-discovery-blue.png"),
    ("A veto is a system event",
     "Logged as system_constraint_veto, class system. It never inflates the "
     "override rate, which is the headline signal for whether the weights are wrong.",
     "icon-checklist-cool-white.png"),
    ("Measured with a permutation test",
     "Reassign each visit randomly among who was actually eligible. High p means "
     "the imbalance comes from eligibility, not scoring: a constraints conversation.",
     "icon-eye-discovery-blue.png"),
]):
    x = M + i * (cw + 0.3)
    rect(s, x, y + 0.9, cw, 2.75, C["coolBlue"], radius=0.2)
    icon(s, ic, x + 0.32, y + 1.15, 0.62)
    text(s, head, x + 0.32, y + 1.92, cw - 0.64, 0.6, size=15.5, bold=True,
         color=C["discovery"], spacing=1.1)
    text(s, body, x + 0.32, y + 2.6, cw - 0.64, 1.0, size=12, color=C["muted"],
         spacing=1.25)
text(s, "On the synthetic pilot: CV 1.06 (RED), permutation p = 0.011. That "
        "imbalance is NOT explained by eligibility, so it needs acting on.",
     M, y + 3.85, W - 2 * M, 0.5, size=13, bold=True, color=C["orange"])

# --- 16. audit log ---------------------------------------------------------
s, y = new("Tracking", "Log the pool, not the pick")
text(s, "Every coordinator considered for every visit gets a row. About five extra "
        "rows per decision, and it is the difference between a system you can "
        "validate and one you can only defend.",
     M, y, W - 2 * M, 0.7, size=15, color=C["ink"], spacing=1.3)
rows = [
    ("scoring_run", "21", "One row per decision. Version fingerprints, pool sizes, "
     "starvation flag, surprise codes"),
    ("candidate_score", "48", "One row per coordinator IN THE POOL. All seven Layer 1 "
     "flags, every raw input, every weighted contribution"),
    ("assignment_outcome", "16", "Who got it, human override vs system veto, reason "
     "class, provisional state, downstream outcomes"),
    ("weight_vector", "20", "Versioned weights with elicitation method, CR, approver"),
    ("calendar_sync_log", "9", "Provider, latency, success, error code. The SLO panel"),
    ("optimizer_shadow", "12", "Weekly greedy vs optimal, regret, escalate flag"),
]
table(s, ["Table", "Cols", "Purpose"], rows, M, y + 0.85, W - 2 * M,
      col_w=[2.5, 0.8, W - 2 * M - 3.3], row_h=0.48, size=12.5,
      align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT])
rect(s, M, y + 3.85, W - 2 * M, 0.8, C["discovery"], radius=0.15)
text(s, "Append-only. No UPDATE, no DELETE. Corrections are new rows, and any past "
        "decision replays under new weights without re-querying a single calendar.",
     M + 0.3, y + 4.03, W - 2 * M - 0.6, 0.5, size=13, bold=True, color=C["coolWhite"])

# --- 17. why the whole pool ------------------------------------------------
s, y = new("Tracking", "What logging the whole pool buys you")
math_img(s, "logit", y + 0.1, x=M, height=1.15)
text(s, "Conditional logit on what schedulers actually chose. Normalised betas are "
        "the weights the humans really use. Needs 50 to 100 decisions with three or "
        "more alternatives.",
     M, y + 1.5, 6.0, 0.9, size=13, color=C["muted"], spacing=1.25)
bullets(s, [
    "Rank-reversal replay under any candidate weight vector",
    "“Why wasn't Kali offered?” answered in one query",
    "Counterfactual re-scoring with no calendar re-query",
    "Criticality and OAT sensitivity on real decisions",
], 6.9, y + 0.15, W - M - 6.9, size=13, gap=0.44)
rect(s, 6.9, y + 2.15, W - M - 6.9, 1.15, C["coolBlue"], radius=0.15)
text(s, "Log only the winner and every one of these becomes permanently "
        "impossible. It costs nothing to do it right on day one.",
     7.2, y + 2.35, W - M - 7.5, 0.8, size=13, bold=True, color=C["discovery"],
     spacing=1.25)
text(s, "Override reasons use a closed vocabulary, split into data_defect (the "
        "inputs were wrong: goes to the fix queue) and preference (the ranking was "
        "right, a human disagreed: goes to weight re-elicitation). Counting them "
        "together is how a scoring system quietly rots.",
     M, y + 3.6, W - 2 * M, 0.9, size=13, color=C["muted"], spacing=1.3)

# --- 18. the debrief -------------------------------------------------------
s, y = new("Weekly debrief", "Unexpected outcomes are detected, not remembered")
rows = [
    ("HUMAN_OVERRODE_TOP", "a human picked past rank 1"),
    ("INSIDE_REVIEW_BAND", "top two inside the calibrated band"),
    ("LOW_SELECTION_STABILITY", "leader wins under 60% of the simplex"),
    ("CRITERION_INERT", "every candidate pinned to the same boundary"),
    ("TOP_PICK_OVER_CAPACITY", "the best option is already over capacity"),
    ("POOL_STARVATION", "one or zero feasible candidates"),
    ("UNEXPLAINED_SCORE_SHIFT", "score moved 0.15 with unchanged inputs"),
]
table(s, ["Detector rule", "Fires when"], rows, M, y, 7.6,
      col_w=[3.7, 3.9], row_h=0.38, size=11.5)
bullets(s, [
    "Eight sections: version header, fairness panel, exception register, "
    "surprise log, override waterfalls, drift panel, decisions, standing question",
    "Waterfalls turn anecdote into data: “suggested C03, chose C06; C03 led on "
    "continuity +0.113, C06 led on burden relief +0.006; net −0.207”",
    "Written to markdown and branded HTML every Monday at 07:00, before the lab "
    "meeting",
], 8.3, y, W - M - 8.3, size=12.5, gap=0.42)
rect(s, M, y + 3.15, W - 2 * M, 0.95, C["discovery"], radius=0.15)
text(s, "Section 8 asks out loud: did any assignment feel wrong without showing up "
        "in sections 2 or 3? That is the only check we have on the detector's false "
        "negatives, and it is why the debrief is a conversation, not an emailed PDF.",
     M + 0.3, y + 3.35, W - 2 * M - 0.6, 0.7, size=13.5, bold=True, color=C["coolWhite"],
     spacing=1.25)

# --- 19. validation plan ---------------------------------------------------
s, y = new("Before the pilot ends", "Three independent routes to the same four numbers")
cw = (W - 2 * M - 0.6) / 3
for i, (head, body, tag) in enumerate([
    ("AHP", "Saaty 1 to 9, six comparisons, individually not by consensus. "
     "Aggregate by the GEOMETRIC MEAN OF JUDGMENTS, which preserves reciprocity. "
     "Accept at CR < 0.10.", "stated preference"),
    ("DEMATEL", "T = N(I-N)^-1. Prominence gives weights, relation gives the "
     "redundancy diagnostic: a strongly negative r-c means the criterion is an "
     "EFFECT and should not carry independent weight.", "interdependence"),
    ("Conditional logit", "Fitted from the override log once 50 to 100 decisions "
     "are in. The weights the schedulers actually use, as opposed to the ones "
     "they say they use.", "revealed preference"),
]):
    x = M + i * (cw + 0.3)
    rect(s, x, y, cw, 2.6, C["coolBlue"], radius=0.2)
    text(s, tag, x + 0.32, y + 0.28, cw - 0.64, 0.25, size=10.5, bold=True,
         color=C["orange"], caps=True, char_spacing=1.2)
    text(s, head, x + 0.32, y + 0.62, cw - 0.64, 0.42, size=19, bold=True,
         color=C["discovery"])
    text(s, body, x + 0.32, y + 1.15, cw - 0.64, 1.3, size=12, color=C["muted"],
         spacing=1.25)
math_img(s, "dematel", y + 2.85, cx=W / 2, height=0.36)
bullets(s, [
    "Sensitivity: one-at-a-time ±0.05 with renormalisation, per-decision "
    "criticality by bisection, Dirichlet Monte Carlo over the simplex",
    "Avoid Chang's extent analysis for fuzzy AHP: it is known to assign zero "
    "weight to non-dominated criteria, which would silently delete a criterion",
    "Timing: run it at the end of pilot week 3",
], M, y + 3.4, W - 2 * M, size=12.5, gap=0.40)

# --- 20. automation --------------------------------------------------------
s, y = new("Running it", "Four scheduled jobs, nothing outside the project folder")
rows = [
    ("calsync", "every 5 min", "Delta pull. Keeps the 72-hour class inside its "
     "15-minute hard threshold with room to spare"),
    ("reconcile", "nightly 02:00", "Full reconcile plus an append-only integrity "
     "check: a shrinking row count means something else touched the database"),
    ("shadow", "Monday 06:45", "Shadow optimiser, so the regret numbers exist "
     "before the debrief renders"),
    ("debrief", "Monday 07:00", "Drift metrics and the debrief report, ahead of "
     "the lab meeting"),
]
table(s, ["Job", "Schedule", "What it does"], rows, M, y, W - 2 * M,
      col_w=[1.7, 2.1, W - 2 * M - 3.8], row_h=0.58, size=12.5)
text(s, "make init  ·  make demo  ·  make test  ·  make week  ·  "
        "make debrief  ·  make install-automation",
     M, y + 2.75, W - 2 * M, 0.4, size=14, bold=True, color=C["discovery"],
     font=FH)
bullets(s, [
    "The debrief job REPORTS a recalibrated review band but never writes it. "
    "Weights and bands change only at a version boundary, with a human approving",
    "Removal leaves data, reports and logs in place: the audit log outlives the "
    "automation",
    "25 correctness anchors run in under ten seconds, including the hand-computed "
    "three-coordinator reference case",
], M, y + 3.3, W - 2 * M, size=12.5, gap=0.42)

# --- 21. asks --------------------------------------------------------------
n += 1
s = slide(C["discovery"])
logo_pair(s, "cool-white")
eyebrow(s, "What we need from you", M, 1.25, C["coolWhite"])
text(s, "Six decisions, none of them long", M, 1.6, 10.5, 0.8, size=34, bold=True,
     color=C["coolWhite"], char_spacing=-0.3)
asks = [
    ("Gamma, in one question", "How many extra minutes of clinic time would you "
     "accept to avoid ten minutes of driving? Team median ÷ 10.", "scheduler, week 1"),
    ("Real capacity per coordinator", "capacity_hours_week from actual FTE, not a "
     "flat 20.", "lab manager, week 1"),
    ("AHP session", "Six pairwise comparisons, twenty minutes, individually.",
     "PI + team, end of week 3"),
    ("Phi vs P at rho = 0.6", "Borderline redundant on synthetic data. If "
     "real data pushes it over, P folds into Phi or becomes a constraint.",
     "PI, week 4"),
    ("Override discipline", "Every override gets a reason code. An unexplained "
     "override is a lost data point.", "everyone, from week 1"),
    ("IRB note", "Randomised tie-breaking as a QI activity, and a retention window "
     "for family-linked scheduling data.", "PI, before go-live"),
]
cw = (W - 2 * M - 0.5) / 2
for i, (head, body, who) in enumerate(asks):
    col, row = i % 2, i // 2
    x = M + col * (cw + 0.5)
    ty = 2.7 + row * 1.42
    rect(s, x, ty, 0.075, 1.15, C["yellow"])
    text(s, head, x + 0.28, ty, cw - 0.4, 0.32, size=16, bold=True, color=C["coolWhite"])
    text(s, body, x + 0.28, ty + 0.36, cw - 0.4, 0.62, size=11.5, color=C["coolWhite"],
         spacing=1.2)
    text(s, who, x + 0.28, ty + 0.98, cw - 0.4, 0.22, size=10, bold=True,
         color=C["yellow"], caps=True, char_spacing=1.0)
footer(s, n, dark=True)


# ---------------------------------------------------------------------------

THEME_COLORS = [
    ("dk1", "000000"), ("lt1", "FFFFFF"), ("dk2", "1A1A1F"), ("lt2", "F4F4F6"),
    ("accent1", "3366FF"), ("accent2", "91BAF4"), ("accent3", "F57F00"),
    ("accent4", "F4DA26"), ("accent5", "D74E2D"), ("accent6", "F8B2B1"),
    ("hlink", "3366FF"), ("folHlink", "91BAF4"),
]


def brand_the_theme(path: str) -> None:
    """Rewrite the stock Office theme with the ESD palette and typeface.

    python-pptx ships the default Office theme, so anyone who later opens this
    deck and picks a colour from the PowerPoint palette gets Office blue rather
    than discovery blue. All the styling here is set at slide level and does not
    depend on the theme, but leaving a wrong palette in the file is an invitation
    to introduce drift on the next edit.
    """
    import re
    import shutil
    import zipfile

    src = zipfile.ZipFile(path, "r")
    items = {n: src.read(n) for n in src.namelist()}
    src.close()

    theme_name = next(n for n in items if n.startswith("ppt/theme/theme"))
    xml = items[theme_name].decode("utf-8")
    for tag, hex_value in THEME_COLORS:
        xml = re.sub(
            rf"<a:{tag}>.*?</a:{tag}>",
            f'<a:{tag}><a:srgbClr val="{hex_value}"/></a:{tag}>',
            xml,
            count=1,
            flags=re.S,
        )
    xml = re.sub(r'(<a:majorFont><a:latin typeface=")[^"]*"',
                 r'\1Libre Franklin"', xml, count=1)
    xml = re.sub(r'(<a:minorFont><a:latin typeface=")[^"]*"',
                 r'\1Libre Franklin Medium"', xml, count=1)
    items[theme_name] = xml.encode("utf-8")

    tmp = path + ".tmp"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as out:
        for name, blob in items.items():
            out.writestr(name, blob)
    shutil.move(tmp, path)


def main() -> None:
    prs.save(OUT)
    brand_the_theme(OUT)
    print(f"wrote {OUT}")
    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
          f"{W} x {H} in")


if __name__ == "__main__":
    main()
