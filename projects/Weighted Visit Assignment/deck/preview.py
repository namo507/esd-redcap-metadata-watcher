#!/usr/bin/env python3
"""Rasterise the built .pptx to PNGs for visual QA.

LibreOffice is not installed here, so this reads the generated file back through
python-pptx and draws it with PIL using the real Libre Franklin TTFs. It is not a
pixel-perfect PowerPoint renderer, but it reads the *actual saved file* rather
than the build script's intentions, which is what makes it useful: text overflow,
collisions, off-slide shapes and missing images all show up.

    python3 deck/preview.py [file.pptx] [--scale 110]
"""

from __future__ import annotations

import os
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

ROOT = os.path.dirname(os.path.abspath(__file__))
OUTDIR = os.path.join(ROOT, "preview")
DEFAULT_ASSETS = (
    "/Users/namomac/Library/Application Support/Claude/local-agent-mode-sessions/"
    "skills-plugin/c898a6b1-6fa4-467a-8bf4-94014b33c02c/"
    "07ab60b6-de84-467a-ae49-82bf671a695e/skills/esd-lab/assets"
)
FONT_DIR = os.path.join(os.environ.get("ESD_ASSETS", DEFAULT_ASSETS), "fonts")

FONTS = {
    ("Libre Franklin", True): "LibreFranklin-Bold.ttf",
    ("Libre Franklin", False): "LibreFranklin-Regular.ttf",
    ("Libre Franklin Medium", True): "LibreFranklin-Bold.ttf",
    ("Libre Franklin Medium", False): "LibreFranklin-Medium.ttf",
}
_cache: dict = {}


def font_for(name: str, bold: bool, px: int):
    key = (name, bold, px)
    if key not in _cache:
        filename = FONTS.get((name, bold), "LibreFranklin-Medium.ttf")
        _cache[key] = ImageFont.truetype(os.path.join(FONT_DIR, filename), px)
    return _cache[key]


def emu_to_px(v, scale):
    return int(round(Emu(v).inches * scale))


def rgb(color, default=(0, 0, 0)):
    try:
        if color and color.rgb is not None:
            r = color.rgb
            return (r[0], r[1], r[2]) if isinstance(r, (bytes, tuple)) else (
                int(str(r)[0:2], 16), int(str(r)[2:4], 16), int(str(r)[4:6], 16)
            )
    except Exception:  # noqa: BLE001
        pass
    return default


def wrap(draw, words, font, max_w):
    lines, cur = [], ""
    for word in words.split(" "):
        trial = f"{cur} {word}".strip()
        if draw.textlength(trial, font=font) <= max_w or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines


def render_slide(slide, w_px, h_px, scale, index):
    img = Image.new("RGB", (w_px, h_px), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    warnings = []

    for shape in slide.shapes:
        try:
            x = emu_to_px(shape.left, scale)
            y = emu_to_px(shape.top, scale)
            w = emu_to_px(shape.width, scale)
            h = emu_to_px(shape.height, scale)
        except (TypeError, AttributeError):
            continue

        if x < -2 or y < -2 or x + w > w_px + 2 or y + h > h_px + 2:
            warnings.append(
                f"shape off-slide: {shape.shape_type} at "
                f"({x/scale:.2f}, {y/scale:.2f}) size {w/scale:.2f}x{h/scale:.2f} in"
            )

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                import io

                pic = Image.open(io.BytesIO(blob)).convert("RGBA")
                pic = pic.resize((max(1, w), max(1, h)), Image.LANCZOS)
                img.paste(pic, (x, y), pic)
            except Exception as exc:  # noqa: BLE001
                warnings.append(f"picture failed: {exc}")
            continue

        if shape.has_text_frame and not shape.text_frame.text.strip():
            fill = None
            try:
                if shape.fill.type is not None and shape.fill.type == 1:
                    fill = rgb(shape.fill.fore_color, (230, 238, 252))
            except Exception:  # noqa: BLE001
                pass
            if fill:
                draw.rounded_rectangle([x, y, x + w, y + h], radius=6, fill=fill)
            continue

        # Autoshape with a fill (our rect() helper) may still carry no text.
        try:
            if shape.fill.type == 1 and not shape.has_text_frame:
                draw.rounded_rectangle(
                    [x, y, x + w, y + h], radius=6, fill=rgb(shape.fill.fore_color)
                )
                continue
        except Exception:  # noqa: BLE001
            pass

        if not shape.has_text_frame:
            continue

        cy = y
        for para in shape.text_frame.paragraphs:
            runs = [r for r in para.runs if r.text]
            if not runs:
                cy += 6
                continue
            run = runs[0]
            size_pt = run.font.size.pt if run.font.size else 14
            px = max(6, int(round(size_pt * scale / 72.0)))
            f = font_for(run.font.name or "Libre Franklin Medium",
                         bool(run.font.bold), px)
            colour = rgb(run.font.color, (26, 26, 31))
            content = "".join(r.text for r in runs)
            line_h = int(px * (para.line_spacing or 1.15) * 1.02)
            for line in wrap(draw, content, f, max(10, w)):
                tw = draw.textlength(line, font=f)
                tx = x
                if para.alignment == PP_ALIGN.CENTER:
                    tx = x + (w - tw) / 2
                elif para.alignment == PP_ALIGN.RIGHT:
                    tx = x + w - tw
                draw.text((tx, cy), line, font=f, fill=colour)
                cy += line_h
        if cy > y + h + max(4, 0.10 * scale):
            warnings.append(
                f"text overflows its box by {(cy - y - h) / scale:.2f} in: "
                f"“{shape.text_frame.text[:52].strip()}…”"
            )
        if cy > h_px:
            warnings.append(f"text runs off the bottom of slide {index}")
    return img, warnings


def main() -> None:
    path = sys.argv[1] if len(sys.argv) > 1 and not sys.argv[1].startswith("-") else (
        os.path.join(os.path.dirname(ROOT), "ESD-Visit-Scheduling-v3.pptx")
    )
    scale = 110
    if "--scale" in sys.argv:
        scale = int(sys.argv[sys.argv.index("--scale") + 1])

    prs = Presentation(path)
    w_px = emu_to_px(prs.slide_width, scale)
    h_px = emu_to_px(prs.slide_height, scale)
    os.makedirs(OUTDIR, exist_ok=True)

    total = 0
    for i, slide in enumerate(prs.slides, start=1):
        img, warnings = render_slide(slide, w_px, h_px, scale, i)
        img.save(os.path.join(OUTDIR, f"slide-{i:02d}.png"))
        for warning in warnings:
            print(f"  slide {i:02d}: {warning}")
            total += 1
    print(f"\n{len(prs.slides._sldIdLst)} slides -> {OUTDIR}  ({total} warning(s))")


if __name__ == "__main__":
    main()
